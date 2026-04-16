"""
Extract text from .pptx files using python-pptx.
Preserves slide structure, speaker notes, and table content.

Usage: python3 extract_pptx.py <file_path>
"""

import sys
import os


def extract_pptx(path):
    """Extract all text content from a .pptx file."""
    from pptx import Presentation

    prs = Presentation(path)
    output = []

    for slide_num, slide in enumerate(prs.slides, 1):
        output.append(f"\n{'='*60}")
        output.append(f"SLIDE {slide_num}")
        output.append(f"{'='*60}")

        # Extract text from all shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        output.append(text)

            # Extract table content
            if shape.has_table:
                table = shape.table
                output.append("\n[Table]")
                for row in table.rows:
                    cells = [cell.text.strip().replace('\n', ' ')
                             for cell in row.cells]
                    output.append(" | ".join(cells))

        # Extract speaker notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                output.append(f"\n[Speaker Notes]")
                output.append(notes)

    return "\n".join(output)


def main():
    if len(sys.argv) != 2:
        print("Usage: python3 extract_pptx.py <file_path>", file=sys.stderr)
        sys.exit(1)

    path = sys.argv[1]
    if not os.path.exists(path):
        print(f"Error: file not found — {path}", file=sys.stderr)
        sys.exit(1)

    ext = os.path.splitext(path)[1].lower()
    if ext != '.pptx':
        print(f"Error: unsupported format '{ext}'. Use .pptx (convert .ppt first)", file=sys.stderr)
        sys.exit(1)

    print(extract_pptx(path))


if __name__ == "__main__":
    main()
