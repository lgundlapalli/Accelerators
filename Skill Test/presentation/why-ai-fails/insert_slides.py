"""
Insert 3 new slides between slides 8 and 9 of Cultural Transformation.pptx.
Sources:
  - Accelerating AI Innovation (Crucial Learning eBook)
  - Cultural Debt / Workforce Transformation deck (April 24, 2026)
"""

import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import shutil
import os

SRC = "Cultural Transformation.pptx"
DST = "Cultural Transformation - Updated.pptx"

shutil.copy(SRC, DST)
prs = Presentation(DST)

# ── helpers ──────────────────────────────────────────────────────────────────

def clone_slide_layout(prs, layout_index=1):
    """Return a blank slide using a layout from the deck."""
    return prs.slide_layouts[layout_index]

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=RGBColor(0x1F,0x1F,0x1F),
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def move_slide_to(prs, from_index, to_index):
    """Move slide at from_index to to_index."""
    xml_slides = prs.slides._sldIdLst
    slides = prs.slides
    # get slide id element
    slide = slides[from_index]
    rId = slides._sldIdLst[from_index].get('r:id')
    # remove from current position
    elem = xml_slides[from_index]
    xml_slides.remove(elem)
    # insert at new position
    xml_slides.insert(to_index, elem)

def duplicate_slide(prs, slide_index):
    """Duplicate a slide and append it to the end. Returns new slide."""
    template_slide = prs.slides[slide_index]
    blank_layout = prs.slide_layouts[6]  # blank
    new_slide = prs.slides.add_slide(blank_layout)
    # Copy the XML content
    new_slide.shapes._spTree.clear()
    for elem in template_slide.shapes._spTree:
        new_slide.shapes._spTree.append(copy.deepcopy(elem))
    # Copy background
    if template_slide.background.fill.type is not None:
        new_slide.background.fill.solid()
        new_slide.background.fill.fore_color.rgb = \
            template_slide.background.fill.fore_color.rgb
    return new_slide


# ── slide dimensions ──────────────────────────────────────────────────────────
W = prs.slide_width   # typically 9144000 EMU (10 in)
H = prs.slide_height  # typically 5143500 EMU (7.5 in)
M = Inches(0.5)       # margin

# Colors from existing deck (navy/dark bg slides inferred)
NAVY   = RGBColor(0x1A, 0x2B, 0x4A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0x7A, 0xC2)  # blue accent
GOLD   = RGBColor(0xF0, 0xA5, 0x00)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF9)
DARK_TXT = RGBColor(0x1F, 0x1F, 0x1F)
GRAY   = RGBColor(0x60, 0x60, 0x60)

# ── use slide 8 as template (index 7) to clone styling ───────────────────────
# We'll build 3 fresh slides and insert them at index 8 (after slide 8)

blank_layout = prs.slide_layouts[6]  # completely blank

# ─────────────────────────────────────────────────────────────────────────────
# NEW SLIDE A: "Behaviorally Agile Organizations Adopt AI 12x Faster"
# Source: Accelerating AI Innovation (Crucial Learning)
# ─────────────────────────────────────────────────────────────────────────────
slideA = prs.slides.add_slide(blank_layout)

# Dark navy background
bg = slideA.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = NAVY

# Title bar accent line
add_rect(slideA, M, Inches(0.3), W - 2*M, Inches(0.06), ACCENT)

# Title
add_textbox(slideA,
    "Behaviorally Agile Organizations Adopt AI 12× Faster",
    M, Inches(0.45), W - 2*M, Inches(0.8),
    font_size=28, bold=True, color=WHITE)

# Subtitle
add_textbox(slideA,
    "The human system — not the technology — determines who wins the AI race.",
    M, Inches(1.2), W - 2*M, Inches(0.5),
    font_size=16, italic=True, color=RGBColor(0xAA, 0xCC, 0xFF))

# Three stat boxes
stats = [
    ("12×", "Faster adoption\nvs. slowest orgs", ACCENT),
    ("6×",  "More likely to achieve\nfull breakthrough adoption", GOLD),
    ("2×",  "More likely to build\na habit of continuous improvement", GREEN),
]

box_w = Inches(2.6)
box_h = Inches(2.8)
box_top = Inches(1.85)
gap = Inches(0.3)
start_left = (W - 3*box_w - 2*gap) / 2

for i, (num, label, color) in enumerate(stats):
    left = start_left + i * (box_w + gap)
    # card bg
    card = slideA.shapes.add_shape(1, left, box_top, box_w, box_h)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x28, 0x3A, 0x5A)
    card.line.color.rgb = color

    # big number
    num_box = slideA.shapes.add_textbox(left, box_top + Inches(0.2), box_w, Inches(1.2))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    run.font.size = Pt(52)
    run.font.bold = True
    run.font.color.rgb = color

    # label
    lbl_box = slideA.shapes.add_textbox(left + Inches(0.15), box_top + Inches(1.3),
                                         box_w - Inches(0.3), Inches(1.3))
    lbl_box.text_frame.word_wrap = True
    lp = lbl_box.text_frame.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run()
    lr.text = label
    lr.font.size = Pt(14)
    lr.font.color.rgb = WHITE

# Source note
add_textbox(slideA,
    "Source: Crucial Learning — Survey of 1,700 professionals across industries",
    M, H - Inches(0.55), W - 2*M, Inches(0.4),
    font_size=10, color=GRAY, italic=True)

# Speaker notes
slideA.notes_slide.notes_text_frame.text = (
    "Lead with the headline stat: 'The research is unambiguous — the fastest organizations adopt new technology twelve times faster than the slowest. That gap isn't explained by better tools or bigger budgets. It's explained by culture.'\n\n"
    "Walk the three numbers. The 6x stat is the most important for this audience: the behaviorally agile orgs aren't just faster — they're six times more likely to achieve FULL adoption. Not pilot success. Full adoption.\n\n"
    "Source: Crucial Learning surveyed 1,700 professionals across every major industry. The pattern held consistently."
)


# ─────────────────────────────────────────────────────────────────────────────
# NEW SLIDE B: "Four Behaviors That Separate Fast Adopters from Slow Ones"
# Source: Accelerating AI Innovation (Crucial Learning)
# ─────────────────────────────────────────────────────────────────────────────
slideB = prs.slides.add_slide(blank_layout)

bg = slideB.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = LIGHT_BG

add_rect(slideB, M, Inches(0.3), W - 2*M, Inches(0.06), ACCENT)

add_textbox(slideB,
    "Four Behaviors That Separate Fast Adopters from Slow Ones",
    M, Inches(0.45), W - 2*M, Inches(0.75),
    font_size=26, bold=True, color=NAVY)

add_textbox(slideB,
    "Behaviorally agile organizations share four measurable communication patterns — absent in slow adopters.",
    M, Inches(1.15), W - 2*M, Inches(0.5),
    font_size=14, color=GRAY)

behaviors = [
    ("1. Speak Up",    "3×", "more likely to raise concerns\nor surface new ideas early",
     "Silence lets fear and uncertainty stall adoption."),
    ("2. Remind",      "5×", "more likely to call out peers\nwho skip new approaches",
     "Peer accountability sustains adoption after launch."),
    ("3. Confront",    "2.5×","more likely to hold resistors\naccountable for non-adoption",
     "Silence in the face of resistance signals the initiative doesn't matter."),
    ("4. Challenge",   "4×", "more likely to challenge assumptions,\neven when leaders are invested",
     "Sacred cows survive when dissent is unsafe."),
]

bw = (W - 2*M - Inches(0.3)) / 2
bh = Inches(1.55)
b_gap_x = Inches(0.3)
b_gap_y = Inches(0.2)
b_top_start = Inches(1.75)

for idx, (title, mult, desc, example) in enumerate(behaviors):
    col = idx % 2
    row = idx // 2
    left = M + col * (bw + b_gap_x)
    top  = b_top_start + row * (bh + b_gap_y)

    card = slideB.shapes.add_shape(1, left, top, bw, bh)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = ACCENT
    card.line.width = Pt(1.5)

    # behavior title + multiplier on same line area
    header_box = slideB.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1),
                                            bw - Inches(0.3), Inches(0.45))
    header_box.text_frame.word_wrap = True
    hp = header_box.text_frame.paragraphs[0]
    hr = hp.add_run()
    hr.text = f"{title}  —  "
    hr.font.bold = True
    hr.font.size = Pt(13)
    hr.font.color.rgb = NAVY
    hr2 = hp.add_run()
    hr2.text = mult
    hr2.font.bold = True
    hr2.font.size = Pt(15)
    hr2.font.color.rgb = ACCENT

    desc_box = slideB.shapes.add_textbox(left + Inches(0.15), top + Inches(0.52),
                                          bw - Inches(0.3), Inches(0.65))
    desc_box.text_frame.word_wrap = True
    dp = desc_box.text_frame.paragraphs[0]
    dr = dp.add_run()
    dr.text = desc
    dr.font.size = Pt(11)
    dr.font.color.rgb = DARK_TXT

    ex_box = slideB.shapes.add_textbox(left + Inches(0.15), top + Inches(1.1),
                                        bw - Inches(0.3), Inches(0.4))
    ex_box.text_frame.word_wrap = True
    ep = ex_box.text_frame.paragraphs[0]
    er = ep.add_run()
    er.text = example
    er.font.size = Pt(10)
    er.font.italic = True
    er.font.color.rgb = GRAY

add_textbox(slideB,
    "Source: Crucial Learning — Accelerating AI Innovation",
    M, H - Inches(0.55), W - 2*M, Inches(0.4),
    font_size=10, color=GRAY, italic=True)

slideB.notes_slide.notes_text_frame.text = (
    "These four behaviors account for almost half the difference between the fastest and slowest cohorts in the study.\n\n"
    "Ask the audience to self-score: 'Think about your own team. How many of these are present today?'\n\n"
    "The key insight is that these aren't personality traits — they're learnable behaviors. Organizations can deliberately build them through training, modeling, and reinforcement.\n\n"
    "Transition: 'But the window to act is narrowing. The workforce disruption is already underway.'"
)


# ─────────────────────────────────────────────────────────────────────────────
# NEW SLIDE C: "The Workforce Disruption Is Already Underway — Act Before the Gap Widens"
# Source: Cultural Debt / Workforce Transformation deck (Gartner 2025)
# ─────────────────────────────────────────────────────────────────────────────
slideC = prs.slides.add_slide(blank_layout)

bg = slideC.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = NAVY

add_rect(slideC, M, Inches(0.3), W - 2*M, Inches(0.06), GOLD)

add_textbox(slideC,
    "The Workforce Disruption Is Already Underway",
    M, Inches(0.45), W - 2*M, Inches(0.75),
    font_size=28, bold=True, color=WHITE)

add_textbox(slideC,
    "Organizations that delay transformation investment won't just fall behind — they'll face compounding talent shortages.",
    M, Inches(1.15), W - 2*M, Inches(0.55),
    font_size=14, italic=True, color=RGBColor(0xAA, 0xCC, 0xFF))

urgency_stats = [
    ("32M",    "Jobs redesigned annually\nthrough 2031",           ACCENT),
    ("39%",    "Of workers' core skills\nchanging by 2030",        GOLD),
    ("2–5 yr", "New half-life of technical\nskills (was 8–12 yrs)", RGBColor(0xFF,0x6B,0x35)),
    ("30%",    "Terminated-by-AI employees\nrehired at higher cost", GREEN),
]

sw = Inches(2.0)
sh = Inches(2.4)
s_top = Inches(1.85)
s_gap = Inches(0.26)
s_start = (W - 4*sw - 3*s_gap) / 2

for i, (num, lbl, color) in enumerate(urgency_stats):
    left = s_start + i * (sw + s_gap)
    card = slideC.shapes.add_shape(1, left, s_top, sw, sh)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x28, 0x3A, 0x5A)
    card.line.color.rgb = color
    card.line.width = Pt(2)

    nb = slideC.shapes.add_textbox(left, s_top + Inches(0.15), sw, Inches(1.0))
    np_ = nb.text_frame.paragraphs[0]
    np_.alignment = PP_ALIGN.CENTER
    nr = np_.add_run()
    nr.text = num
    nr.font.size = Pt(36)
    nr.font.bold = True
    nr.font.color.rgb = color

    lb = slideC.shapes.add_textbox(left + Inches(0.1), s_top + Inches(1.1),
                                    sw - Inches(0.2), Inches(1.2))
    lb.text_frame.word_wrap = True
    lp_ = lb.text_frame.paragraphs[0]
    lp_.alignment = PP_ALIGN.CENTER
    lr_ = lp_.add_run()
    lr_.text = lbl
    lr_.font.size = Pt(12)
    lr_.font.color.rgb = WHITE

# Bottom call-to-action bar
add_rect(slideC, M, Inches(4.45), W - 2*M, Inches(0.7), RGBColor(0x0D, 0x1E, 0x38))

add_textbox(slideC,
    "The CIO mandate: build the human system before competitors do — or spend 3× more to fix it later.",
    M + Inches(0.2), Inches(4.5), W - 2*M - Inches(0.4), Inches(0.6),
    font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(slideC,
    "Source: Gartner Strategic Planning Assumptions (2025); WEF Future of Jobs Report (2025)",
    M, H - Inches(0.55), W - 2*M, Inches(0.4),
    font_size=10, color=GRAY, italic=True)

slideC.notes_slide.notes_text_frame.text = (
    "This slide connects the cultural transformation argument to hard workforce economics.\n\n"
    "Walk the numbers: '32 million jobs are being redesigned EVERY YEAR through 2031 — that's not a future scenario, it's happening now. 39% of your employees' core skills will need to change by 2030. The half-life of a technical skill has collapsed from 8-12 years to 2-5 years.'\n\n"
    "The 30% rehire stat is the most concrete business case: 'Organizations that cut roles to replace with AI are spending the next 18 months hiring those people back — often at 30-50% higher cost — because they didn't invest in the transition. That's a cost of delay, not a cost of transformation.'\n\n"
    "Close with the bottom bar: 'This is the math. Addressing cultural debt early costs a leadership conversation. Addressing it late costs a transformation program.'\n\n"
    "Source: Gartner Strategic Planning Assumptions (2025); World Economic Forum Future of Jobs Report (2025)."
)


# ─────────────────────────────────────────────────────────────────────────────
# Move the 3 new slides (currently at end) to positions 8, 9, 10 (0-indexed)
# so they sit between original slide 8 (idx 7) and original slide 9 (idx 8)
# After adding 3 slides, deck has N+3 slides.
# New slides are at indices N, N+1, N+2.
# We want them at indices 8, 9, 10.
# ─────────────────────────────────────────────────────────────────────────────

total = len(prs.slides)  # N+3
# Move slide at (total-3) → position 8, then (total-2) → 9, then (total-1) → 10
# We need to move in order, accounting for shifting indices.

xml_slides = prs.slides._sldIdLst

def move_to(xml_slides, from_idx, to_idx):
    elem = xml_slides[from_idx]
    xml_slides.remove(elem)
    xml_slides.insert(to_idx, elem)

# Slides A, B, C are at positions total-3, total-2, total-1
# Move A to index 8 first
move_to(xml_slides, total - 3, 8)
# B is now at total-2 (shifted by 1 since we inserted before it)
move_to(xml_slides, total - 2, 9)
# C is now at total-1
move_to(xml_slides, total - 1, 10)

prs.save(DST)
print(f"Saved: {DST}")
print(f"Total slides: {len(prs.slides)}")
