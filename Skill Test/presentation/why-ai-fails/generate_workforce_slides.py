"""
Standalone deck: Workforce Transformation Imperative (new slides only)
Abbott BTS-DTS branded — April 24, 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x00, 0x00, 0x50)
LAVENDER = RGBColor(0xDD, 0xDD, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
MED_BLUE = RGBColor(0x33, 0x33, 0xAA)
LIGHT_GRAY = RGBColor(0x88, 0x88, 0x88)
SOFT_LAVENDER = RGBColor(0xF0, 0xF0, 0xFF)
LIGHT_BLUE_ACCENT = RGBColor(0xBB, 0xBB, 0xFF)
LIGHT_NAVY = RGBColor(0x99, 0x99, 0xCC)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

page_counter = [0]


def next_page():
    page_counter[0] += 1
    return page_counter[0]


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_footer(slide, page_num=None):
    if page_num is None:
        page_num = page_counter[0]
    tf = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(8), Inches(0.4)).text_frame
    p = tf.paragraphs[0]
    p.text = "Proprietary and confidential — do not distribute"
    p.font.size = Pt(9)
    p.font.color.rgb = LIGHT_GRAY
    p.font.italic = True
    tf2 = slide.shapes.add_textbox(Inches(12), Inches(6.9), Inches(1), Inches(0.4)).text_frame
    p2 = tf2.paragraphs[0]
    p2.text = str(page_num)
    p2.font.size = Pt(9)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.RIGHT


def navy_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY


def add_text(slide, left, top, width, height, text, font_size=18, bold=False,
             color=NEAR_BLACK, alignment=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.font.italic = italic
    p.font.name = 'Helvetica Neue'
    return tf


def add_box(slide, left, top, width, height, fill_color):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.fill.background()
    return box


def add_table(slide, left, top, width, height, rows, cols, data,
              col_widths=None, font_size=11, header_color=LAVENDER):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = 'Helvetica Neue'
                paragraph.font.color.rgb = NEAR_BLACK
                paragraph.space_after = Pt(2)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = NAVY
            else:
                cell.fill.background()
    return table


# ============================================================
# SLIDE 1: Title (Navy)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
navy_bg(slide)
add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.2),
         "The Workforce Transformation Imperative", font_size=44, color=WHITE)
add_text(slide, Inches(0.8), Inches(2.9), Inches(11), Inches(0.8),
         "AI Is Reshaping Work — How CIOs Must Respond",
         font_size=28, color=WHITE)
add_text(slide, Inches(0.8), Inches(4.0), Inches(10), Inches(0.6),
         "Based on Gartner Strategic Planning Assumptions (2025)",
         font_size=16, italic=True, color=LIGHT_BLUE_ACCENT)
add_text(slide, Inches(0.8), Inches(5.2), Inches(10), Inches(0.4),
         "BTS-DTS Architecture  |  April 24, 2026  |  v1.0", font_size=14, color=LIGHT_NAVY)

add_notes(slide, "Open by grounding the room: 'We've been talking about AI and jobs for years. The question now isn't whether AI will change work — it already is. The question is whether your organization is positioned to lead that change or scramble to catch up.' Don't read the slide. Set the stakes.")


# ============================================================
# SLIDE 2: The Scale of Change (Big Numbers)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
pg = next_page()
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "AI Is Reshaping Work — Not Eliminating It — But Reinvention Is Urgent",
         font_size=24, color=NAVY)
add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.4),
         "The AI-driven jobs apocalypse narrative is misplaced. The real risk is failing to reinvent.",
         font_size=13, italic=True, color=NEAR_BLACK)

data = [
    ["The Reality", "The Shift", "The Risk"],
    ["32M", "60%", "30%"],
    ["Jobs requiring refactoring every year through 2031",
     "Of digital products architected for AI agents by 2029, with human UX secondary",
     "Of employees terminated by AI will be rehired at higher cost by 2029"],
    ["AI is not eliminating jobs — it is redesigning them at unprecedented scale",
     "Autonomous agents initiate, negotiate, and complete tasks without human input — requiring a new design paradigm centered on agent experience",
     "Ineffective workforce transition strategies create the very talent shortages they were meant to avoid"],
    ["39%", "2–5 years", "15%"],
    ["Of workers' core skills expected to change by 2030 (WEF Future of Jobs 2025)",
     "New half-life of technical skills by 2030 — down from 8–12 years",
     "Lower employee engagement when AI agents appear in org charts alongside humans"],
]
add_table(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(4.5), 6, 3, data,
          col_widths=[Inches(4), Inches(4), Inches(4)], font_size=11)

add_text(slide, Inches(0.6), Inches(6.2), Inches(12), Inches(0.3),
         "Sources: Gartner (2025); World Economic Forum Future of Jobs Report (2025)",
         font_size=9, italic=True, color=LIGHT_GRAY)
add_footer(slide, pg)

add_notes(slide, """Start by saying the obvious out loud: 'Everyone is worried AI is going to eliminate jobs. The data says that's not quite right — but what is happening is actually harder to manage.' Then walk the numbers left to right. Thirty-two million jobs a year aren't disappearing — they're being redesigned. That means the people doing those jobs need different skills, different tools, and different support. The half-life number is the one to land on: we used to expect a technical skill to stay relevant for 8 to 12 years. Now it's 2 to 5. That means if you hired someone for their Python or data skills in 2023, some of that is already degrading. The 30% rehire stat is the one that surprises people most — organizations cut roles, then realize six months later they needed those people, and end up paying more to bring them back or replace them. The point is: this isn't a story about job loss, it's a story about organizations being unprepared for the speed of change.""")


# ============================================================
# SLIDE 3: Five Strategic Planning Assumptions
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
pg = next_page()
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "Five Strategic Predictions That Will Redefine the Workforce by 2030",
         font_size=26, color=NAVY)

data = [
    ["By When", "What Changes", "Why It Matters"],
    ["2030",
     "Half-life of technical skills shrinks to 2–5 years (from 8–12)",
     "Adaptability and learning velocity become the primary hiring metric — not credentials"],
    ["2029",
     "60% of digital products architected primarily for AI agent consumption",
     "Human-facing UX becomes secondary; CIOs must evolve toward agent experience design"],
    ["2029",
     "30% of employees terminated and replaced by AI will be rehired — often at higher cost",
     "Ineffective workforce transitions create acute skill shortages and drive up acquisition cost"],
    ["2028",
     "At least one large enterprise will claim the right to create AI avatars of all employees",
     "Employment contracts must address digital identity, knowledge replication, and consent"],
    ["2028",
     "Organizations displaying AI agents in team structures will have 15% lower engagement",
     "AI agents belong in governed IAM/HR systems — not in org charts alongside employees"],
]
add_table(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(5.0), 6, 3, data,
          col_widths=[Inches(1.5), Inches(5.25), Inches(5.25)], font_size=11)

add_text(slide, Inches(0.6), Inches(6.3), Inches(12), Inches(0.3),
         "Source: Gartner Strategic Planning Assumptions (2025)",
         font_size=9, italic=True, color=LIGHT_GRAY)
add_footer(slide, pg)

add_notes(slide, """Walk through each row like a timeline. By 2030, the skills that get people hired today will have a shelf life of 2 to 5 years — so the question you ask in interviews shifts from 'what do you know?' to 'how fast do you learn?' By 2029, most digital products will be built primarily for AI agents to use, not humans — that's a complete flip from how we design today. The rehire stat comes next: companies that cut people to replace them with AI will spend the next few years hiring them back, often at higher cost, because they didn't plan the transition well. The avatar row is the one that generates the most conversation — some companies already want to create digital replicas of employees using their emails, documents, and decisions. The legal frameworks don't exist yet. And the last one is practical: if you put AI agents in your org chart next to people, employees disengage. Keep AI governance in IT systems, not in reporting structures. For each row, the useful question is: 'Are we ready for this, and if not, what would it take?'""")


# ============================================================
# SLIDE 4: Four CIO Actions
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
pg = next_page()
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "Four Actions to Future-Proof the Workforce",
         font_size=28, color=NAVY)

actions = [
    ("Invest in adaptive learning systems",
     "Personalize learning programs, proactively recommend microlearning courses, verify skills in real time — replace annual training with continuous capability development.",
     "CHRO + CIO"),
    ("Redesign for human-agent collaboration",
     "Train UX and engineering teams to build workflows where human judgment integrates with AI-driven actions — and that can eventually support machine-to-machine interactions.",
     "CTO + UX Leadership"),
    ("Architect the talent remix",
     "Partner with HR to redefine roles, skills, and workforce design to realize AI ambitions — any workforce restructuring decision must flow from this talent strategy, not precede it.",
     "CIO + CHRO"),
    ("Lead on digital identity governance",
     "Drive enterprisewide dialogue on responsible use of employee data in AI and digital avatars — stress-test employment contracts now for scenarios where AI replicates human skill.",
     "CIO + General Counsel + HR"),
]

for i, (title, desc, owner) in enumerate(actions):
    y = Inches(1.1) + Inches(i * 1.3)
    color = LAVENDER if i % 2 == 0 else SOFT_LAVENDER
    add_box(slide, Inches(0.6), y, Inches(12), Inches(1.15), color)
    add_text(slide, Inches(1.0), y + Inches(0.05), Inches(9), Inches(0.35),
             title, font_size=15, bold=True, color=NAVY)
    add_text(slide, Inches(1.0), y + Inches(0.4), Inches(10.5), Inches(0.45),
             desc, font_size=11, color=NEAR_BLACK)
    add_text(slide, Inches(1.0), y + Inches(0.85), Inches(10), Inches(0.25),
             f"Leads: {owner}", font_size=10, bold=True, color=MED_BLUE)

add_box(slide, Inches(0.6), Inches(6.35), Inches(12), Inches(0.45), NAVY)
add_text(slide, Inches(1.0), Inches(6.4), Inches(11), Inches(0.35),
         "The CIO’s mandate: create a symbiotic relationship between humans and machines before competitors do.",
         font_size=13, bold=True, color=WHITE, italic=True)

add_text(slide, Inches(0.6), Inches(6.85), Inches(12), Inches(0.25),
         "Source: Gartner (2025)",
         font_size=9, italic=True, color=LIGHT_GRAY)
add_footer(slide, pg)

add_notes(slide, """There are four things on this slide and none of them are IT problems alone — each one needs IT and HR working together, or IT and Legal. Walk through them in order. First, adaptive learning: the old model was send people to a class once a year. That doesn't work anymore. You need systems that notice when someone's skills are drifting and push them the right content before it becomes a problem. Second, human-agent collaboration design: your UX and engineering teams are building products for how humans and AI work together today — but they also need to think about what happens when the AI is mostly talking to other AI systems. That's a design shift. Third, the talent remix: before you make any workforce restructuring decision, you need a clear picture of what roles you actually need in an AI-first world. Otherwise you cut the wrong people or create gaps you can't fill. Fourth, digital identity governance: if someone in your org is already building an AI model trained on employee data, do your employment contracts cover that? Probably not. Get ahead of it now. The broader point is that the CIO can't wait for HR to lead this — if you do, you'll already be behind by the time the conversation starts.""")


# ============================================================
# SAVE
# ============================================================
output = "/Users/GUNDLLX/learn-claude/Skill Test/presentation/why-ai-fails/cultural debt april 24 2026.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Slides: {len(prs.slides)}")
