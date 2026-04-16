"""
'The Hidden Balance Sheet: Why AI Initiatives Fail and What Cultural Debt Costs You'
Abbott BTS-DTS branded PPTX generator.
Full deck structure: Title > Exec Summary > Agenda > Section Dividers + Content > Appendix
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
NAVY = RGBColor(0x00, 0x00, 0x50)
LAVENDER = RGBColor(0xDD, 0xDD, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
MED_BLUE = RGBColor(0x33, 0x33, 0xAA)
LIGHT_GRAY = RGBColor(0x88, 0x88, 0x88)
SOFT_LAVENDER = RGBColor(0xF0, 0xF0, 0xFF)
RED_ACCENT = RGBColor(0xCC, 0x33, 0x33)
GREEN_ACCENT = RGBColor(0x22, 0x88, 0x22)
LIGHT_BLUE_ACCENT = RGBColor(0xBB, 0xBB, 0xFF)
LIGHT_NAVY = RGBColor(0x99, 0x99, 0xCC)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

page_counter = [0]  # mutable counter for auto-incrementing


def next_page():
    page_counter[0] += 1
    return page_counter[0]


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_footer(slide, page_num=None):
    if page_num is None:
        page_num = page_counter[0]
    tf = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(8), Inches(0.4)).text_frame
    p = tf.paragraphs[0]
    p.text = "Proprietary and confidential \u2014 do not distribute"
    p.font.size = Pt(9)
    p.font.color.rgb = LIGHT_GRAY
    p.font.italic = True
    tf2 = slide.shapes.add_textbox(Inches(12), Inches(6.9), Inches(1), Inches(0.4)).text_frame
    p2 = tf2.paragraphs[0]
    p2.text = str(page_num)
    p2.font.size = Pt(9)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.RIGHT


def navy_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY


def add_text(slide, left, top, width, height, text, font_size=18, bold=False,
             color=NEAR_BLACK, alignment=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.font.italic = italic
    p.font.name = 'Helvetica Neue'
    return tf


def add_box(slide, left, top, width, height, fill_color):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.fill.background()
    return box


def add_table(slide, left, top, width, height, rows, cols, data,
              col_widths=None, font_size=11, header_color=LAVENDER):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = 'Helvetica Neue'
                paragraph.font.color.rgb = NEAR_BLACK
                paragraph.space_after = Pt(2)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = NAVY
            else:
                cell.fill.background()
    return table


def add_section_divider(number, title, subtitle=None):
    """Add a navy section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    navy_bg(slide)
    add_text(slide, Inches(0.8), Inches(2.5), Inches(4), Inches(0.8),
             number, font_size=48, color=WHITE)
    add_text(slide, Inches(0.8), Inches(3.5), Inches(10), Inches(1),
             title, font_size=40, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.8), Inches(4.5), Inches(10), Inches(0.6),
                 subtitle, font_size=20, italic=True, color=LIGHT_BLUE_ACCENT)
    return slide


# ============================================================
# SLIDE 1: Title (Navy Background)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
navy_bg(slide)
add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.2),
         "The Hidden Balance Sheet", font_size=48, color=WHITE)
add_text(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(1),
         "Why AI Initiatives Fail and What Cultural Debt Costs You",
         font_size=32, color=WHITE)
add_text(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(0.6),
         "A case for treating culture as infrastructure \u2014 not afterthought",
         font_size=16, italic=True, color=LIGHT_BLUE_ACCENT)
add_text(slide, Inches(0.8), Inches(5.5), Inches(10), Inches(0.4),
         "BTS-DTS Architecture  |  April 2026", font_size=14, color=LIGHT_NAVY)

add_notes(slide, """OPENING (15 seconds):
"We've invested heavily in AI. Today I want to talk about why most AI initiatives \u2014 across the industry, not just here \u2014 aren't delivering returns. And more importantly, what we can do about it."

Don't read the slide. Set the room with a direct statement.

TRANSITION:
"Let me give you the executive summary, then we'll go section by section."
""")


# ============================================================
# SLIDE 2: Executive Summary (3-Column Grid)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
pg = next_page()
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "Executive Summary", font_size=32, color=NAVY)
add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.5),
         "AI initiatives fail at 70-85% \u2014 not from bad technology, but from cultural debt that compounds with every failed attempt.",
         font_size=14, color=NEAR_BLACK)

data = [
    ["The Problem", "The Hidden Cost", "The Path Forward"],
    ["70-85%", "3x", "Culture before code"],
    ["AI initiative failure rate across industries. Technology is rarely the cause.",
     "Higher failure rate on second attempt when cultural debt is unaddressed.",
     "Organizations that invest in cultural readiness see 2-3x higher adoption."],
    ["What's failing", "What's compounding", "What works"],
    ["Not the models, not the data, not the platforms \u2014 the organization itself.",
     "Each failure breeds cynicism, shadow processes, and passive resistance.",
     "Psychological safety, aligned incentives, and visible leadership commitment."],
]
add_table(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(4.2), 5, 3, data,
          col_widths=[Inches(4), Inches(4), Inches(4)], font_size=11)

add_text(slide, Inches(0.6), Inches(6.1), Inches(12), Inches(0.4),
         "Sources: McKinsey Global Institute (2023); Gartner AI in the Enterprise (2024); MIT CISR (2024)",
         font_size=9, italic=True, color=LIGHT_GRAY)
add_footer(slide, pg)

add_notes(slide, """PURPOSE: This is the "read only this slide" slide. If someone reads ONLY this, they get the full story.

KEY TALKING POINTS:
- "Three columns: the problem is clear \u2014 70-85% failure rate. The hidden cost is that it compounds \u2014 each failure makes the next one harder. And the path forward is counterintuitive: invest in culture before technology."
- Emphasize the 3x stat: "If your first AI initiative fails and you don't address the cultural root cause, your second attempt faces roughly triple the resistance."

TRANSITION:
- "Here's what we'll cover."
""")


# ============================================================
# SLIDE 3: Agenda (Numbered Sections)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
pg = next_page()
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "What We'll Explore Today", font_size=32, color=NAVY)

agenda = [
    ("01", "The Numbers", "Why 70-85% of AI initiatives fail \u2014 and technology isn't the reason"),
    ("02", "The Root Cause", "What cultural debt is, how it compounds, and why no one tracks it"),
    ("03", "The Diagnosis", "Five cultural patterns that predict AI failure before launch"),
    ("04", "The Path Forward", "What organizations that break the cycle do differently"),
    ("05", "The Actions", "Three things you can do this quarter to start paying down cultural debt"),
]
for i, (num, title, desc) in enumerate(agenda):
    y = Inches(1.4) + Inches(i * 1.0)
    add_text(slide, Inches(1.2), y, Inches(1.2), Inches(0.6),
             num, font_size=28, bold=True, color=NAVY)
    add_text(slide, Inches(2.8), y + Inches(0.0), Inches(4), Inches(0.4),
             title, font_size=18, bold=True, color=NAVY)
    add_text(slide, Inches(2.8), y + Inches(0.35), Inches(9), Inches(0.4),
             desc, font_size=13, color=NEAR_BLACK)
    if i < len(agenda) - 1:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(1.2), y + Inches(0.8), Inches(10.5), Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        line.line.fill.background()

add_footer(slide, pg)

add_notes(slide, """Don't spend more than 15 seconds here. Simply say:
"Five sections \u2014 we start with the data, define the root cause, diagnose where we're stuck, show what works, and end with three concrete actions."

Move quickly to the numbers.
""")


# ============================================================
# SLIDE 4: Section Divider \u2014 01 The Numbers
# ============================================================
slide = add_section_divider("01", "The Numbers")
add_notes(slide, """Pause briefly. Let the room reset before the data slide.
"Let's start with what the research tells us."
""")


# ============================================================
# SLIDE 5: Big Numbers \u2014 The Failure Rate
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
pg = next_page()
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "70-85% of AI Initiatives Fail \u2014 Technology Isn\u2019t the Reason",
         font_size=28, color=NAVY)
add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.5),
         "The technology works. The organization doesn\u2019t change.",
         font_size=14, color=NEAR_BLACK)

metrics = [
    ("70-85%", "Fail to deliver\nexpected value", "Consistent across industries,\ncompany sizes, and geographies"),
    ("<15%", "Failures attributed\nto technology", "Models work. Infrastructure works.\nThe organization doesn\u2019t adopt."),
    ("#1 Cause", "Organizational and\ncultural resistance", "Change resistance, skill gaps,\nmisaligned incentives, lack of trust"),
    ("$4.6T", "Projected global AI\ninvestment by 2027", "The majority at risk of\ndelivering no return"),
]
for i, (number, label, desc) in enumerate(metrics):
    x = Inches(0.6) + Inches(i * 3.1)
    add_text(slide, x, Inches(1.8), Inches(2.8), Inches(1),
             number, font_size=44, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(2.9), Inches(2.8), Inches(0.7),
             label, font_size=13, bold=True, color=MED_BLUE, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.7), Inches(2.8), Inches(1),
             desc, font_size=11, color=NEAR_BLACK, alignment=PP_ALIGN.CENTER)

add_box(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(0.7), LAVENDER)
add_text(slide, Inches(1), Inches(5.1), Inches(11), Inches(0.5),
         "The gap between AI capability and organizational readiness is where value dies.",
         font_size=16, bold=True, color=NAVY)

add_text(slide, Inches(0.6), Inches(6.2), Inches(12), Inches(0.3),
         "Sources: McKinsey (2023); BCG AI Adoption Survey (2024); IDC AI Spending Guide (2025)",
         font_size=9, italic=True, color=LIGHT_GRAY)
add_footer(slide, pg)

add_notes(slide, """DELIVERY: Walk left to right through each number. Let each one land.

70-85%: "Across every major survey \u2014 McKinsey, Gartner, BCG \u2014 the failure rate is consistent."

<15%: "Less than 15% of failures are attributed to technology. The models work. The organization doesn\u2019t adopt."

#1 Cause: "The number one reason: organizational and cultural resistance."

$4.6T: "Global AI investment is projected at $4.6 trillion by 2027. The majority at risk."

Point to the callout bar. Pause.

TRANSITION: "So if it\u2019s not the technology, what is it?"
""")


# ============================================================
# SLIDE 6: Section Divider \u2014 02 The Root Cause
# ============================================================
slide = add_section_divider("02", "The Root Cause", "What cultural debt is and how it compounds")
add_notes(slide, """"This is the concept that explains why the numbers look the way they do." """)


# ============================================================
# SLIDE 7: Cultural Debt vs Technical Debt (Side-by-Side)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
pg = next_page()
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Cultural Debt Compounds Like Technical Debt \u2014 But No One Tracks It",
         font_size=26, color=NAVY)

# Left panel: Technical Debt
add_box(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(3.8), SOFT_LAVENDER)
add_text(slide, Inches(0.9), Inches(1.4), Inches(5.2), Inches(0.4),
         "Technical Debt", font_size=20, bold=True, color=NAVY)

tech_items = [
    "Shortcuts in code that create future rework",
    "Visible in code reviews, build times, bug rates",
    "Tracked, measured, budgeted for",
    "Compounds as system complexity grows",
    "Fix: refactor the code",
]
for i, item in enumerate(tech_items):
    add_text(slide, Inches(0.9), Inches(2.0) + Inches(i * 0.5), Inches(5.2), Inches(0.45),
             f"\u2022  {item}", font_size=13, color=NEAR_BLACK)

# Right panel: Cultural Debt
add_box(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(3.8), LAVENDER)
add_text(slide, Inches(7.1), Inches(1.4), Inches(5.2), Inches(0.4),
         "Cultural Debt", font_size=20, bold=True, color=NAVY)

culture_items = [
    "Shortcuts in change management that create future resistance",
    "Invisible until the next initiative stalls",
    "Untracked, unmeasured, unbudgeted",
    "Compounds as organizational cynicism grows",
    "Fix: rebuild trust, realign incentives, invest in people",
]
for i, item in enumerate(culture_items):
    add_text(slide, Inches(7.1), Inches(2.0) + Inches(i * 0.5), Inches(5.2), Inches(0.45),
             f"\u2022  {item}", font_size=13, color=NEAR_BLACK)

# Bottom insight
add_box(slide, Inches(0.6), Inches(5.4), Inches(12), Inches(0.8), NAVY)
add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
         "Cultural debt accrues every time an organization deploys technology without investing in the humans who must adopt it.",
         font_size=15, bold=True, color=WHITE, italic=True)

add_footer(slide, pg)

add_notes(slide, """This is the definitional slide \u2014 it creates the vocabulary the rest of the deck builds on.

FRAMING:
"Everyone in this room understands technical debt. But there\u2019s a parallel concept that\u2019s far more dangerous because it\u2019s invisible: cultural debt."

Walk the comparison. Pause on row 3: "Tracked, measured, budgeted \u2014 versus untracked, unmeasured, unbudgeted."

Read the bottom bar aloud. Let it land.

TRANSITION: "And here\u2019s why this is urgent \u2014 cultural debt compounds."
""")


# ============================================================
# SLIDE 8: The Compounding Cycle
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
pg = next_page()
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Every Failed AI Initiative Makes the Next One Harder",
         font_size=28, color=NAVY)

# Cycle visualization: 6 boxes in a loop
cycle_items = [
    (Inches(1.5), Inches(1.8), "AI Initiative\nLaunched"),
    (Inches(5.0), Inches(1.8), "Adoption\nStalls"),
    (Inches(8.5), Inches(1.8), "Leadership\nDeclares Failure"),
    (Inches(8.5), Inches(3.8), "Cynicism\nDeepens"),
    (Inches(5.0), Inches(3.8), "Shadow Processes\nHarden"),
    (Inches(1.5), Inches(3.8), "Next Initiative Faces\nHigher Resistance"),
]
for x, y, text in cycle_items:
    add_box(slide, x, y, Inches(2.8), Inches(1.2), LAVENDER)
    add_text(slide, x + Inches(0.1), y + Inches(0.15), Inches(2.6), Inches(0.9),
             text, font_size=14, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)

# Arrows between boxes
arrows = [
    (Inches(4.3), Inches(2.2), "\u2192"),
    (Inches(7.8), Inches(2.2), "\u2192"),
    (Inches(10.2), Inches(3.2), "\u2193"),
    (Inches(7.8), Inches(4.2), "\u2190"),
    (Inches(4.3), Inches(4.2), "\u2190"),
    (Inches(1.8), Inches(3.2), "\u2191"),
]
for x, y, arrow in arrows:
    add_text(slide, x, y, Inches(0.6), Inches(0.5),
             arrow, font_size=28, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)

# Cycle stages at bottom
add_box(slide, Inches(0.6), Inches(5.3), Inches(12), Inches(1.3), SOFT_LAVENDER)
stages = [
    ("Cycle 1:", "Healthy skepticism \u2014 \u201cLet\u2019s see if this works\u201d"),
    ("Cycle 2:", "Active resistance \u2014 \u201cWe tried this before\u201d"),
    ("Cycle 3:", "Organizational antibodies \u2014 \u201cAI doesn\u2019t work here\u201d"),
    ("Cycle 4:", "Structural inability \u2014 the best people have stopped engaging"),
]
for i, (label, desc) in enumerate(stages):
    y = Inches(5.4) + Inches(i * 0.28)
    add_text(slide, Inches(1), y, Inches(1.5), Inches(0.3),
             label, font_size=11, bold=True, color=NAVY)
    add_text(slide, Inches(2.5), y, Inches(9.5), Inches(0.3),
             desc, font_size=11, color=NEAR_BLACK)

add_footer(slide, pg)

add_notes(slide, """Walk through the cycle slowly. Each loop is worse than the last.

"An AI initiative launches. Adoption stalls. Leadership declares failure. Cynicism deepens. Shadow processes harden. And the next initiative faces higher resistance before it even starts."

Then walk the four cycles:
- Cycle 1: Healthy skepticism. Fine.
- Cycle 2: Active resistance. Fighting institutional memory.
- Cycle 3: Organizational antibodies. Resistance is unconscious.
- Cycle 4: Structural inability. The best people have disengaged.

TRANSITION: "So what are the specific patterns that predict this?"
""")


# ============================================================
# SLIDE 9: Section Divider \u2014 03 The Diagnosis
# ============================================================
slide = add_section_divider("03", "The Diagnosis", "Five patterns that predict failure before launch")
add_notes(slide, """"Now we get practical. These five patterns are diagnostic \u2014 if even two are present, the initiative is at risk." """)


# ============================================================
# SLIDE 10: Five Failure Patterns
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
pg = next_page()
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Five Cultural Patterns Predict AI Failure Before Launch",
         font_size=28, color=NAVY)

data = [
    ["Pattern", "What It Looks Like", "What It Costs"],
    ["Permission Culture",
     "Teams wait for approval before experimenting",
     "Innovation dies in committee; competitors move faster"],
    ["Metric Misalignment",
     "AI success measured by deployment, not adoption",
     "Teams ship models no one uses; declare victory anyway"],
    ["Invisible Middle",
     "Middle management neither sponsors nor blocks \u2014 they wait",
     "Initiatives lose momentum in the layer that controls execution"],
    ["Skills Theater",
     "Training programs check boxes but don\u2019t build capability",
     "Teams attend workshops, return to old workflows Monday morning"],
    ["Trust Deficit",
     "\u201cAI will replace us\u201d narrative goes unaddressed",
     "Best talent disengages or leaves; adoption becomes coercion"],
]
add_table(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(4.5), 6, 3, data,
          col_widths=[Inches(2.5), Inches(4.75), Inches(4.75)], font_size=12)

add_box(slide, Inches(0.6), Inches(5.9), Inches(12), Inches(0.7), LAVENDER)
add_text(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.5),
         "If even two of these patterns are present, the initiative is at risk before a single model is trained.",
         font_size=14, bold=True, color=NAVY)

add_footer(slide, pg)

add_notes(slide, """Walk down each row:

"Permission Culture \u2014 teams need approvals before experimenting. Innovation dies in committee."

"Metric Misalignment \u2014 the most common. We measure 'models deployed' not 'decisions improved.'"

"Invisible Middle \u2014 the least discussed, most lethal. Middle managers aren\u2019t resisting \u2014 they\u2019re rationally waiting."

"Skills Theater \u2014 everyone went to training. Nobody changed how they work."

"Trust Deficit \u2014 'AI will replace us' goes unaddressed. Best people quietly disengage."

Ask: "Mentally score your organization. How many are present?"

TRANSITION: "The organizations that break this cycle have three things in common."
""")


# ============================================================
# SLIDE 11: Section Divider \u2014 04 The Path Forward
# ============================================================
slide = add_section_divider("04", "The Path Forward", "What organizations that succeed do differently")
add_notes(slide, """"We\u2019ve seen the problem. Now the turn: what actually works." """)


# ============================================================
# SLIDE 12: What Works
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
pg = next_page()
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Organizations That Break the Cycle Invest in Culture Before Code",
         font_size=26, color=NAVY)

data = [
    ["What They Do", "Why It Works", "Evidence"],
    ["Psychological safety first \u2014\ncreate space for experimentation",
     "People adopt what they help build; fear\nof failure kills adoption faster than bad tech",
     "Google Project Aristotle: psychological\nsafety is #1 predictor of team effectiveness"],
    ["Aligned incentives \u2014 measure\nadoption and impact, not deployment",
     "What gets measured gets done; if the metric\nis \u201cmodel deployed\u201d nobody owns adoption",
     "Microsoft AI transformation: shifted KPIs\nfrom deployment to outcome \u2014 adoption +40%"],
    ["Visible leadership commitment \u2014\nleaders use the tools publicly",
     "Permission flows downward; if the VP doesn\u2019t\nuse it, the team won\u2019t either",
     "MIT CISR: executive engagement is the\nstrongest predictor of transformation success"],
]
add_table(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(3.8), 4, 3, data,
          col_widths=[Inches(4), Inches(4), Inches(4)], font_size=11)

add_box(slide, Inches(0.6), Inches(5.3), Inches(12), Inches(1.2), NAVY)
add_text(slide, Inches(1), Inches(5.4), Inches(11), Inches(0.4),
         "The common thread:", font_size=13, color=LIGHT_NAVY)
add_text(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
         "They treat culture as infrastructure \u2014 something you build and maintain \u2014 not as a side effect of technology deployment.",
         font_size=16, bold=True, color=WHITE)

add_footer(slide, pg)

add_notes(slide, """The turn from problem to solution.

"Psychological safety sounds soft until you realize it\u2019s the difference between teams who experiment and teams who hide."

"Aligned incentives sounds obvious until you audit how many AI programs measure 'models in production' rather than 'decisions improved.'"

"Visible leadership \u2014 not sponsorship in a steering committee, but a VP who shows the team how they used Claude this week."

Read the bottom bar: "They treat culture as infrastructure."

TRANSITION: "So what does the shift look like?"
""")


# ============================================================
# SLIDE 13: The Shift \u2014 People, Process, Technology
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
pg = next_page()
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "The Shift Requires People, Process, and Technology \u2014 In That Order",
         font_size=24, color=NAVY)
add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.4),
         "Most organizations invest in reverse order. That sequence is the source of cultural debt.",
         font_size=14, bold=True, color=MED_BLUE)

# Three columns
cols_data = [
    ("People", Inches(0.6), [
        ("\u2022 From: AI as threat \u2192 To: AI as multiplier", True),
        ("  Invest in capability, not just training. Redefine roles around Human + AI.", False),
        ("\u2022 From: Skills theater \u2192 To: Apprenticeship", True),
        ("  Learn by doing, not by slide deck. Build internal multipliers.", False),
    ]),
    ("Process", Inches(4.6), [
        ("\u2022 From: Deploy then adopt \u2192 To: Co-design then deploy", True),
        ("  Embed adoption into delivery. Measure outcomes, not outputs.", False),
        ("\u2022 From: IT-led \u2192 To: Business-led, IT-enabled", True),
        ("  Business owns the outcome. Technology teams enable.", False),
    ]),
    ("Technology", Inches(8.6), [
        ("\u2022 From: Build it, they will come \u2192 To: Build what they asked for", True),
        ("  Meet people where they are. Integrate into existing workflows.", False),
        ("\u2022 From: Centralized tools \u2192 To: Contextual guardrails", True),
        ("  Responsible AI at platform level. Teams focus on value.", False),
    ]),
]

for title, x, items in cols_data:
    add_box(slide, x, Inches(1.5), Inches(3.8), Inches(4.5), LAVENDER)
    add_text(slide, x + Inches(0.2), Inches(1.6), Inches(3.4), Inches(0.4),
             title, font_size=18, bold=True, color=NAVY)
    y_offset = Inches(2.15)
    for text, is_bold in items:
        fs = 10 if is_bold else 10
        add_text(slide, x + Inches(0.2), y_offset, Inches(3.4), Inches(0.45),
                 text, font_size=fs, bold=is_bold, color=NAVY if is_bold else NEAR_BLACK)
        y_offset += Inches(0.45) if is_bold else Inches(0.55)

# Bottom: Transition Plan
add_box(slide, Inches(0.6), Inches(6.2), Inches(12), Inches(0.5), NAVY)
add_text(slide, Inches(1), Inches(6.25), Inches(11), Inches(0.4),
         "Start with People (Q1)  |  Redesign Process around adoption (Q2)  |  Then select Technology (Q3)",
         font_size=13, bold=True, color=WHITE)

add_footer(slide, pg)

add_notes(slide, """This slide reframes the investment sequence.

"Almost every failed AI initiative started with technology selection. The ones that succeed start with people."

The From/To framing lets the audience self-diagnose.

Read the bottom bar: the sequence matters. People first, process, then technology.

TRANSITION: "And there\u2019s urgency \u2014 because cultural debt doesn\u2019t wait."
""")


# ============================================================
# SLIDE 14: Section Divider \u2014 05 The Actions
# ============================================================
slide = add_section_divider("05", "The Actions", "Three things you can do this quarter")
add_notes(slide, """"We\u2019ve diagnosed the problem. Here\u2019s what to do about it." """)


# ============================================================
# SLIDE 15: Urgency \u2014 Compounding Cost
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
pg = next_page()
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Cultural Debt Grows Quarterly \u2014 Every Delay Compounds the Cost",
         font_size=28, color=NAVY)

# Timeline boxes with escalating colors
stages = [
    ("Quarter 1", "Skepticism", "\u201cLet\u2019s wait and see\u201d", "Low cost to address", LAVENDER),
    ("Quarter 2", "Resistance", "\u201cThis didn\u2019t work last time\u201d", "Moderate cost", SOFT_LAVENDER),
    ("Quarter 4", "Antibodies", "\u201cAI isn\u2019t for us\u201d", "High cost \u2014 requires visible reset", RGBColor(0xFF, 0xE0, 0xB2)),
    ("Year 2", "Structural Lock-in", "The best people stopped engaging", "Requires org redesign", RGBColor(0xFF, 0xCC, 0xBC)),
]

for i, (period, stage, quote, cost, color) in enumerate(stages):
    x = Inches(0.6) + Inches(i * 3.1)
    add_box(slide, x, Inches(1.5), Inches(2.8), Inches(3.5), color)
    add_text(slide, x + Inches(0.2), Inches(1.6), Inches(2.4), Inches(0.4),
             period, font_size=14, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(2.1), Inches(2.4), Inches(0.5),
             stage, font_size=22, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(2.8), Inches(2.4), Inches(0.5),
             quote, font_size=12, italic=True, color=NEAR_BLACK, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(3.6), Inches(2.4), Inches(0.5),
             cost, font_size=11, bold=True, color=MED_BLUE, alignment=PP_ALIGN.CENTER)

# Arrows between stages
for i in range(3):
    x = Inches(3.4) + Inches(i * 3.1)
    add_text(slide, x, Inches(2.5), Inches(0.6), Inches(0.5),
             "\u2192", font_size=32, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)

# Bottom insight
add_box(slide, Inches(0.6), Inches(5.4), Inches(12), Inches(1.0), NAVY)
add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.35),
         "The math is simple:", font_size=14, color=LIGHT_NAVY)
add_text(slide, Inches(1), Inches(5.9), Inches(11), Inches(0.4),
         "Addressing cultural debt in Q1 costs a leadership conversation. Addressing it in Year 2 costs a transformation program.",
         font_size=16, bold=True, color=WHITE)

add_footer(slide, pg)

add_notes(slide, """The urgency slide. Compounding metaphor works because leadership thinks in financial terms.

"Every quarter unaddressed, the cost roughly doubles \u2014 not because the problem gets harder, but because the cynicism gets deeper."

Walk left to right. Colors darken \u2014 visual cue for escalating severity.

TRANSITION: "So here are three actions we can take today."
""")


# ============================================================
# SLIDE 16: Three Actions \u2014 Call to Action
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
pg = next_page()
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Three Actions to Start Paying Down Cultural Debt Today",
         font_size=28, color=NAVY)

actions = [
    ("1. Audit the Cultural Balance Sheet",
     "Run a cultural readiness assessment before the next AI investment. Name the debt.",
     "Survey team sentiment on AI. Map the last 3 failed initiatives. Identify failure patterns.",
     "Executive Sponsor + HR  |  Within 30 days"),
    ("2. Realign the Scorecard",
     "Change how you measure AI success from \u201cdeployed\u201d to \u201cadopted and delivering value.\u201d",
     "Rewrite KPIs for the current AI portfolio. Add adoption metrics. Remove vanity metrics.",
     "Business Unit Leads  |  Within 30 days"),
    ("3. Make Leadership Visible",
     "Executives use the tools, share what they learned, and admit what didn\u2019t work.",
     "One executive shares an AI experiment (success or failure) at the next all-hands. Publicly.",
     "C-Suite  |  Next all-hands"),
]

for i, (title, desc, detail, owner) in enumerate(actions):
    y = Inches(1.2) + Inches(i * 1.65)
    color = LAVENDER if i % 2 == 0 else SOFT_LAVENDER
    add_box(slide, Inches(0.6), y, Inches(12), Inches(1.45), color)
    add_text(slide, Inches(1), y + Inches(0.05), Inches(11), Inches(0.35),
             title, font_size=16, bold=True, color=NAVY)
    add_text(slide, Inches(1), y + Inches(0.4), Inches(11), Inches(0.35),
             desc, font_size=12, color=NEAR_BLACK)
    add_text(slide, Inches(1), y + Inches(0.75), Inches(11), Inches(0.35),
             detail, font_size=11, color=NEAR_BLACK)
    add_text(slide, Inches(1), y + Inches(1.1), Inches(11), Inches(0.3),
             owner, font_size=10, bold=True, color=MED_BLUE)

# Closing line
add_box(slide, Inches(0.6), Inches(6.2), Inches(12), Inches(0.55), NAVY)
add_text(slide, Inches(1), Inches(6.25), Inches(11), Inches(0.45),
         "The question is not whether you\u2019ve invested in AI. The question is whether you\u2019ve invested in the human system required to make it work.",
         font_size=13, bold=True, color=WHITE, italic=True)

add_footer(slide, pg)

add_notes(slide, """End with action, not inspiration. Low-cost, high-signal.

"I\u2019m not asking for a new program or a new budget. Three targeted actions."

1. Cultural audit \u2014 "Name the debt. Can\u2019t fix what you don\u2019t measure."
2. Scorecard \u2014 "Change the metric, change the behavior."
3. Visible leadership \u2014 "Until the C-suite models it, everything else is policy without practice."

Read the closing line. Pause. Let it land.

Ask: "Who in this room owns the first one?"
""")


# ============================================================
# SLIDE 17: Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
pg = next_page()
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "Next Steps", font_size=32, color=NAVY)

data = [
    ["Action", "Owner", "By When"],
    ["Circulate this deck and the cultural debt framework", "[Presenter]", "This week"],
    ["Run cultural readiness assessment on current AI portfolio", "[Executive Sponsor] + HR", "Within 30 days"],
    ["Audit AI program KPIs \u2014 flag any measuring deployment without adoption", "Business Unit Leads", "Within 30 days"],
    ["Schedule leadership AI \u201cshow and tell\u201d", "[C-Suite Champion]", "Next all-hands"],
    ["Reconvene to review findings and agree on debt paydown plan", "[Presenter] + Steering Committee", "45 days"],
]
add_table(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(3.5), 6, 3, data,
          col_widths=[Inches(6.5), Inches(3.5), Inches(2)], font_size=12)

add_footer(slide, pg)

add_notes(slide, """Read the actions, not the table. The 30-day window is important.

If the next meeting isn\u2019t on the calendar before people leave the room, it won\u2019t happen.

Ask: "Can we get the readiness assessment scheduled before we leave today?"
""")


# ============================================================
# SLIDE 18: Appendix Divider (Navy)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
navy_bg(slide)
add_text(slide, Inches(0.8), Inches(3), Inches(10), Inches(1.2),
         "Appendix", font_size=44, color=WHITE)

add_notes(slide, """If within time: "That\u2019s the core presentation. Backup material follows \u2014 research sources, assessment framework, and anticipated questions."

If over time: skip to audience Q&A.""")


# ============================================================
# SLIDE 19: Research Sources
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
pg = next_page()
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "The Research Behind the 70-85% Failure Rate", font_size=28, color=NAVY)

data = [
    ["Source", "Finding", "Year"],
    ["McKinsey Global Institute", "70% of companies report AI transformations fail to reach stated goals", "2023"],
    ["Gartner", "85% of AI projects deliver erroneous outcomes due to bias in data, algorithms, or teams", "2024"],
    ["BCG & MIT Sloan", "Only 10% of companies report significant financial benefit from AI", "2024"],
    ["Rand Corporation", "AI projects fail primarily due to organizational factors, not technology", "2024"],
    ["Deloitte State of AI", "\"Cultural readiness\" as top AI barrier increased from 32% to 47% YoY", "2025"],
]
add_table(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(3.5), 6, 3, data,
          col_widths=[Inches(3.5), Inches(6.5), Inches(2)], font_size=12)

add_footer(slide, pg)

add_notes(slide, """Backup for the "where did you get that number?" question. The 70-85% range uses different study definitions of failure, but the pattern is consistent.""")


# ============================================================
# SLIDE 20: Cultural Debt Assessment Framework
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
pg = next_page()
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "Cultural Debt Assessment Framework", font_size=28, color=NAVY)
add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
         "Rate your organization 1-5 on each dimension. A score below 15 signals significant cultural debt.",
         font_size=13, color=NEAR_BLACK, italic=True)

data = [
    ["Dimension", "1 (High Debt)", "5 (Low Debt)"],
    ["Experimentation Safety", "Failures are punished or hidden", "Failures are shared and learned from openly"],
    ["Incentive Alignment", "AI success = model deployed", "AI success = business outcome improved"],
    ["Middle Management Buy-in", "Managers wait for direction", "Managers actively sponsor experiments"],
    ["Skill Building", "Training is annual, classroom-based", "Learning is continuous, embedded in work"],
    ["Trust Level", "\u201cAI will replace us\u201d is dominant narrative", "\u201cAI will make us better\u201d is believed and evidenced"],
    ["Leadership Visibility", "Leaders talk about AI but don\u2019t use it", "Leaders demo their own AI usage regularly"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(3.5), 7, 3, data,
          col_widths=[Inches(3), Inches(4.5), Inches(4.5)], font_size=11)

# Scoring guide
scoring = [
    ("25-30:", "Cultural infrastructure is strong \u2014 invest in technology with confidence", GREEN_ACCENT),
    ("18-24:", "Moderate debt \u2014 address gaps before scaling AI", MED_BLUE),
    ("12-17:", "Significant debt \u2014 pause technology investment, invest in people", RGBColor(0xCC, 0x88, 0x00)),
    ("6-11:", "Critical debt \u2014 organizational redesign needed before AI can succeed", RED_ACCENT),
]
for i, (score, desc, color) in enumerate(scoring):
    y = Inches(5.2) + Inches(i * 0.35)
    add_text(slide, Inches(1), y, Inches(1.5), Inches(0.3),
             score, font_size=11, bold=True, color=color)
    add_text(slide, Inches(2.5), y, Inches(9.5), Inches(0.3),
             desc, font_size=11, color=NEAR_BLACK)

add_footer(slide, pg)

add_notes(slide, """Most actionable appendix slide. Offer to facilitate the assessment.

The value is in the conversation, not the score. When a VP rates "Experimentation Safety" as 2 and their peer rates it 4, that disagreement is more valuable than the number.""")


# ============================================================
# SLIDE 21: Anticipated Questions
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
pg = next_page()
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "Anticipated Questions", font_size=28, color=NAVY)

questions = [
    ("\u201cIsn\u2019t this just change management?\u201d",
     "Change management addresses a specific initiative. Cultural debt is the accumulated residue of every initiative. If change management were sufficient, the failure rate wouldn\u2019t be 70-85%."),
    ("\u201cWe can\u2019t afford to slow down.\u201d",
     "Every initiative launched into a high-debt org has a 70-85% chance of failing and making the next one harder. One quarter of assessment is cheaper than failing a third time."),
    ("\u201cHow do we measure cultural debt?\u201d",
     "Track adoption rates, not deployment counts. Survey sentiment quarterly. Measure time-to-adoption, not time-to-production. Use the assessment framework in the appendix."),
    ("\u201cWhose budget does this come from?\u201d",
     "The assessment costs almost nothing (survey + facilitated conversation). Interventions use existing budgets redirected, not new spend."),
    ("\u201cWe have executive sponsorship.\u201d",
     "Sponsorship that lives in a steering committee but doesn\u2019t reach the teams doing the work is cultural debt in disguise. The question is whether it\u2019s visible and active."),
]

for i, (q, a) in enumerate(questions):
    y = Inches(1.0) + Inches(i * 1.1)
    add_box(slide, Inches(0.6), y, Inches(5), Inches(0.9), LAVENDER if i % 2 == 0 else SOFT_LAVENDER)
    add_text(slide, Inches(0.8), y + Inches(0.15), Inches(4.6), Inches(0.6),
             q, font_size=12, bold=True, color=NAVY)
    add_text(slide, Inches(5.8), y + Inches(0.05), Inches(7), Inches(0.8),
             a, font_size=11, color=NEAR_BLACK)

add_footer(slide, pg)

add_notes(slide, """Don\u2019t present this \u2014 flip to it when asked.

"Change management" \u2192 comes from HR. Validate, then distinguish.
"Can\u2019t slow down" \u2192 from the CTO. Reframe speed as "deliver value faster."
"Budget" \u2192 easiest: costs almost nothing vs. a failed AI program.""")


# ============================================================
# SAVE
# ============================================================
output = "/Users/GUNDLLX/learn-claude/Skill Test/presentation/why-ai-fails/cultural-debt-presentation.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Slides: {len(prs.slides)}")
