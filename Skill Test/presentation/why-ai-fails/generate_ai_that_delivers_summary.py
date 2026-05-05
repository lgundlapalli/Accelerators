"""
Summary deck: "AI That Delivers" — Crucial Learning (Joseph Grenny & JaeLynn Williams)
Source: Ai Transformation crucial conversatins - Slides.pdf
10 slides capturing the full narrative arc.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height
M = Inches(0.55)

# ── palette (Crucial Learning: black / white / red) ───────────────────────────
BLACK   = RGBColor(0x10, 0x10, 0x10)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
RED     = RGBColor(0xCC, 0x00, 0x00)
LGRAY   = RGBColor(0xF2, 0xF2, 0xF2)
MGRAY   = RGBColor(0xD0, 0xD0, 0xD0)
DGRAY   = RGBColor(0x55, 0x55, 0x55)
OFFWHITE= RGBColor(0xFA, 0xFA, 0xFA)

blank = prs.slide_layouts[6]

# ── helpers ───────────────────────────────────────────────────────────────────

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color, line_color=None, line_pt=0):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h,
       size=18, bold=False, italic=False,
       color=BLACK, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = wrap
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box

def multiline_tb(slide, lines, l, t, w, h,
                 size=16, bold=False, color=BLACK,
                 align=PP_ALIGN.LEFT, line_spacing=1.15):
    """lines = list of (text, bold, color) tuples or plain strings."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, b, c = item, bold, color
        else:
            txt, b, c = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = b
        r.font.color.rgb = c
    return box

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank)
bg(s1, BLACK)

# Red accent stripe top
rect(s1, 0, 0, W, Inches(0.12), RED)

# Title
tb(s1, "AI THAT DELIVERS",
   M, Inches(1.6), W - 2*M, Inches(1.4),
   size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Subtitle
tb(s1, "ALIGNING TECHNOLOGY WITH HUMAN BEHAVIOR",
   M, Inches(3.0), W - 2*M, Inches(0.7),
   size=20, bold=True, color=RED, align=PP_ALIGN.LEFT)

# Divider
rect(s1, M, Inches(3.85), Inches(6), Inches(0.04), WHITE)

# Presenters
tb(s1, "Joseph Grenny & JaeLynn Williams  |  Crucial Learning",
   M, Inches(4.05), Inches(7), Inches(0.5),
   size=16, color=MGRAY)

# Source note
tb(s1, "Source: AI Transformation — Crucial Conversations (April 2026)",
   M, H - Inches(0.7), W - 2*M, Inches(0.45),
   size=11, italic=True, color=DGRAY)

notes(s1, "Title slide. Presenter: Joseph Grenny (Author & Cofounder, Crucial Learning) and JaeLynn Williams (Advisor & Former Executive).")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — The Problem: 95% of AI Pilots Are Failing
# ─────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(blank)
bg(s2, WHITE)

rect(s2, 0, 0, W, Inches(0.1), RED)

tb(s2, "95% of Generative AI Pilots Are Failing",
   M, Inches(0.3), W - 2*M, Inches(0.9),
   size=34, bold=True, color=BLACK)

tb(s2, "The technology works. The human system doesn't change.",
   M, Inches(1.15), W - 2*M, Inches(0.55),
   size=18, italic=True, color=DGRAY)

rect(s2, M, Inches(1.85), Inches(5.8), Inches(3.8), LGRAY, MGRAY, 1)

multiline_tb(s2, [
    ("\"The model was right.  The outcome was wrong.\"", True, BLACK),
    ("", False, BLACK),
    ("Air Methods deployed an AI routing model to reduce crew drive time.", False, DGRAY),
    ("Crews kept driving to their assigned base — 75 miles away —", False, DGRAY),
    ("bypassing a nearby base 8 miles from home.", False, DGRAY),
    ("", False, BLACK),
    ("The model flagged a better route. Nobody acted on it.", False, DGRAY),
    ("The conversation never happened.", False, DGRAY),
],
   M + Inches(0.25), Inches(2.0), Inches(5.3), Inches(3.6),
   size=15, color=BLACK)

# Right side stat
rect(s2, Inches(7.2), Inches(1.85), Inches(5.5), Inches(1.6), RED)
tb(s2, "95%", Inches(7.2), Inches(1.95), Inches(5.5), Inches(0.9),
   size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s2, "of GenAI pilots at companies are failing",
   Inches(7.2), Inches(2.85), Inches(5.5), Inches(0.5),
   size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s2, "— MIT / Fortune (Aug 2025)",
   Inches(7.2), Inches(3.3), Inches(5.5), Inches(0.35),
   size=11, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(s2, Inches(7.2), Inches(3.65), Inches(5.5), Inches(2.0), LGRAY)
multiline_tb(s2, [
    ("AI is deployed as two different things:", True, BLACK),
    ("", False, BLACK),
    ("• Productivity tool  (email, slides, summaries)", False, DGRAY),
    ("• Business transformation tool  (decisions, predictions, redesigned workflows)", False, DGRAY),
    ("", False, BLACK),
    ("The gap between the two is behavioral — not technical.", True, RED),
],
   Inches(7.4), Inches(3.75), Inches(5.1), Inches(1.85),
   size=13, color=BLACK)

notes(s2, "Open with the headline: '95% of GenAI pilots are failing — not because the models are wrong, but because the human system around them doesn't change.' The Air Methods story is the opening case study: a routing AI that correctly identified a closer base, but crews never adopted it. One conversation could have fixed it.")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Root Cause: Behavioral Agility Is the Bottleneck
# ─────────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank)
bg(s3, BLACK)

rect(s3, 0, 0, W, Inches(0.1), RED)

# Left quote block
tb(s3, "BEHAVIORAL\nAGILITY",
   M, Inches(0.3), Inches(5), Inches(1.5),
   size=38, bold=True, color=WHITE)

tb(s3,
   "The health of a relationship, team, or organization\nis a function of the average lag time between\nidentifying and discussing problems.",
   M, Inches(2.0), Inches(5.8), Inches(2.2),
   size=20, italic=True, color=WHITE)

tb(s3,
   "Organizations that act fast on new ideas consistently outperform those that don't.",
   M, Inches(4.2), Inches(5.8), Inches(1.0),
   size=16, color=MGRAY)

# Right side: stat card
rect(s3, Inches(7.0), Inches(0.5), Inches(5.8), Inches(6.2), RGBColor(0x1A,0x1A,0x1A))

tb(s3, "Study of 1,700+ professionals",
   Inches(7.2), Inches(0.7), Inches(5.4), Inches(0.45),
   size=13, italic=True, color=MGRAY, align=PP_ALIGN.CENTER)

stats = [
    ("12×", "faster to implement great ideas"),
    ("6×",  "more consistent in adopting new methods"),
    ("3×",  "more engaged at work"),
]
for i, (num, lbl) in enumerate(stats):
    top = Inches(1.3) + i * Inches(1.7)
    tb(s3, num, Inches(7.2), top, Inches(2.2), Inches(1.1),
       size=52, bold=True, color=RED, align=PP_ALIGN.CENTER)
    tb(s3, lbl, Inches(9.3), top + Inches(0.3), Inches(3.3), Inches(0.7),
       size=15, color=WHITE)

tb(s3, "Behavioral agility is the bottleneck.",
   Inches(7.2), Inches(6.15), Inches(5.4), Inches(0.45),
   size=14, bold=True, italic=True, color=RED, align=PP_ALIGN.CENTER)

notes(s3, "The research finding: the most behaviorally agile teams are 12x faster, 6x more consistent, and 3x more engaged than the least agile. This is the organizational capability gap that determines who wins the AI race.")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Even When AI Is Right, Humans Squander It
# ─────────────────────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(blank)
bg(s4, WHITE)

rect(s4, 0, 0, W, Inches(0.1), RED)

tb(s4, "Even When AI Does Its Job, Humans Can Squander the Opportunity",
   M, Inches(0.25), W - 2*M, Inches(0.95),
   size=30, bold=True, color=BLACK)

tb(s4, "Example: AI for Loan & Claim Decisions",
   M, Inches(1.15), W - 2*M, Inches(0.45),
   size=16, italic=True, color=RED)

# Scenario box
rect(s4, M, Inches(1.65), Inches(4.5), Inches(1.3), LGRAY)
tb(s4, "AI recommends approving a loan applicant who, under previous criteria, would have been denied.",
   M + Inches(0.2), Inches(1.75), Inches(4.1), Inches(1.1),
   size=14, color=BLACK)

# 4 failure modes
problems = [
    ("1  Skepticism",     "Busy underwriter delays action — 'let me think about it'"),
    ("2  Intuition Override", "Past experience says risky → overrides AI 'to be safe'"),
    ("3  Manager Pushback", "Hesitation because manager punishes when AI-based loans go wrong"),
    ("4  Downstream Rejection", "Approves the loan, but the next step in the process declines it"),
]

bw = (W - 2*M - Inches(0.3)) / 2
bh = Inches(2.0)
b_gap = Inches(0.3)
b_top = Inches(3.2)

for idx, (title, desc) in enumerate(problems):
    col = idx % 2
    row = idx // 2
    l = M + col * (bw + b_gap)
    t = b_top + row * (bh + Inches(0.15))
    rect(s4, l, t, bw, bh, LGRAY, RED, 1.5)
    rect(s4, l, t, bw, Inches(0.42), RED)
    tb(s4, title, l + Inches(0.15), t + Inches(0.05), bw - Inches(0.3), Inches(0.35),
       size=13, bold=True, color=WHITE)
    tb(s4, desc, l + Inches(0.15), t + Inches(0.5), bw - Inches(0.3), Inches(1.35),
       size=13, color=BLACK)

notes(s4, "Use the loan underwriting example to make human failure modes concrete. Ask the audience: 'Which of these four happens in your organization when AI delivers a recommendation that contradicts past practice?'")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — The Central Question
# ─────────────────────────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(blank)
bg(s5, BLACK)

rect(s5, 0, 0, W, Inches(0.1), RED)

rect(s5, M, Inches(1.2), Inches(3.5), Inches(0.55), RED)
tb(s5, "OUR QUESTION", M + Inches(0.15), Inches(1.27), Inches(3.2), Inches(0.42),
   size=14, bold=True, color=WHITE)

tb(s5,
   "Why are so few organizations\ntranslating AI business insights\ninto substantial business results?",
   M, Inches(2.1), W - 2*M, Inches(2.5),
   size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(s5, M, Inches(4.85), W - 2*M, Inches(0.04), RED)

tb(s5,
   "AI exposes how your organization actually works.\nThe answer to the question above is the answer to your transformation.",
   M, Inches(5.0), W - 2*M, Inches(0.9),
   size=18, italic=True, color=MGRAY, align=PP_ALIGN.CENTER)

notes(s5, "This is the pivot slide. Let the question sit. Then reframe: 'This isn't a technology question. It's a leadership question.' AI exposes the organizational habits — silence, avoidance, workarounds — that were already there.")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — 4 Crucial Conversations for Behavioral Agility
# ─────────────────────────────────────────────────────────────────────────────
s6 = prs.slides.add_slide(blank)
bg(s6, WHITE)

rect(s6, 0, 0, W, Inches(0.1), RED)

tb(s6, "Four Crucial Conversations That Drive AI Adoption",
   M, Inches(0.25), W - 2*M, Inches(0.85),
   size=32, bold=True, color=BLACK)

tb(s6, "Behaviorally agile organizations have these conversations consistently. Others don't.",
   M, Inches(1.05), W - 2*M, Inches(0.5),
   size=16, italic=True, color=DGRAY)

convos = [
    ("SPEAKING UP",
     "Openly name concerns\nthat signal a need for\ninnovation or course correction.",
     "3× more likely\nto speak up quickly"),
    ("REMINDING",
     "Nudge colleagues who\ndefault to old routines\nwhen new methods are called for.",
     "5× more likely\nto remind peers"),
    ("HOLDING ACCOUNTABLE",
     "Direct, immediate correction\nwhen someone violates agreed\nbehavioral norms.",
     "6× more likely to suggest\nimprovements without\nfear of power dynamics"),
    ("CHALLENGING\nSACRED COWS",
     "Follow the data — even when\nit contradicts leadership\npreferences or past practice.",
     "3× more likely peers\nimmediately correct\nmisuse of new approaches"),
]

cw = (W - 2*M - Inches(0.45)) / 4
ch = Inches(4.6)
c_gap = Inches(0.15)
c_top = Inches(1.7)

for i, (title, desc, stat) in enumerate(convos):
    l = M + i * (cw + c_gap)
    rect(s6, l, c_top, cw, ch, LGRAY)
    rect(s6, l, c_top, cw, Inches(0.7), BLACK)
    tb(s6, title, l + Inches(0.1), c_top + Inches(0.1), cw - Inches(0.2), Inches(0.55),
       size=12, bold=True, color=WHITE)
    tb(s6, desc,
       l + Inches(0.1), c_top + Inches(0.8), cw - Inches(0.2), Inches(1.8),
       size=13, color=BLACK)
    rect(s6, l, c_top + Inches(2.7), cw, Inches(1.7), RED)
    tb(s6, stat,
       l + Inches(0.1), c_top + Inches(2.85), cw - Inches(0.2), Inches(1.4),
       size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

notes(s6, "Walk each column. The stats are from Crucial Learning's study of 1,700 professionals. Key point: these are learnable behaviors — not personality traits. Organizations can deliberately build them through modeling, coaching, and accountability.")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — AI Fails in Predictable Ways (4 Failure Modes)
# ─────────────────────────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(blank)
bg(s7, BLACK)

rect(s7, 0, 0, W, Inches(0.1), RED)

tb(s7, "AI Fails in Predictable Ways",
   M, Inches(0.25), W - 2*M, Inches(0.85),
   size=36, bold=True, color=WHITE)

tb(s7, "These weren't isolated incidents. The same four patterns appear across every industry.",
   M, Inches(1.05), W - 2*M, Inches(0.5),
   size=16, italic=True, color=MGRAY)

failures = [
    ("VISIBILITY",      "Speak Up",    "3M Health Info Systems",
     "ML model outputs were inconsistent. Clinicians changed input forms — revenue cycle never knew. Cross-team silence made the AI unreliable.",
     "This was a visibility failure."),
    ("ADOPTION",        "Remind",      "3M Health Info Systems",
     "ICD-10 expanded codes 100,000x. Coders kept working from memory, then entered into the system after. The rollout took a day. The behavior change took months.",
     "This was an adoption failure."),
    ("ACCOUNTABILITY",  "Confront",    "3M Health Info Systems",
     "Early results were uneven. Customers stayed because they trusted transparent communication — same data, same room, same story across all stakeholders.",
     "This was an accountability failure."),
    ("ASSUMPTIONS",     "Challenge",   "Air Methods",
     "Routing model built on 'strong systems data.' Area managers dismissed crew resistance as change-aversion. One direct conversation surfaced the fatal flaw.",
     "This was an assumption failure."),
]

fw = (W - 2*M - Inches(0.45)) / 4
fh = Inches(5.1)
f_gap = Inches(0.15)
f_top = Inches(1.7)

for i, (dim, behavior, source, story, label) in enumerate(failures):
    l = M + i * (fw + f_gap)
    rect(s7, l, f_top, fw, fh, RGBColor(0x1A, 0x1A, 0x1A))
    rect(s7, l, f_top, fw, Inches(0.5), RED)
    tb(s7, dim, l + Inches(0.1), f_top + Inches(0.07), fw - Inches(0.2), Inches(0.38),
       size=12, bold=True, color=WHITE)
    tb(s7, f"Behavior: {behavior}", l + Inches(0.1), f_top + Inches(0.58), fw - Inches(0.2), Inches(0.4),
       size=11, bold=True, color=RED)
    tb(s7, source, l + Inches(0.1), f_top + Inches(0.95), fw - Inches(0.2), Inches(0.35),
       size=10, italic=True, color=MGRAY)
    tb(s7, story, l + Inches(0.1), f_top + Inches(1.35), fw - Inches(0.2), Inches(2.4),
       size=11, color=WHITE)
    tb(s7, label, l + Inches(0.1), f_top + Inches(3.85), fw - Inches(0.2), Inches(0.45),
       size=11, bold=True, italic=True, color=RED)

notes(s7, "Each of these case studies is from 3M Health Information Systems (except the last, Air Methods). Walk across: 'In every case the model worked. The failure was human.' Use the failure type label to anchor: visibility, adoption, accountability, assumptions.")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — AI Execution Diagnostic™
# ─────────────────────────────────────────────────────────────────────────────
s8 = prs.slides.add_slide(blank)
bg(s8, WHITE)

rect(s8, 0, 0, W, Inches(0.1), RED)

tb(s8, "AI EXECUTION DIAGNOSTIC™",
   M, Inches(0.25), W - 2*M, Inches(0.75),
   size=36, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
tb(s8, "Score your AI initiative in 60 seconds  |  0 = Red  |  1 = Mixed  |  2 = Green  (per dimension, 0–8 total)",
   M, Inches(0.95), W - 2*M, Inches(0.45),
   size=13, italic=True, color=DGRAY, align=PP_ALIGN.CENTER)

dims = [
    ("VISIBILITY",      "What don't we see?",
     "Fragmented inputs,\nhidden dependencies",
     "Same inputs,\nshared visibility"),
    ("ADOPTION",        "Are people using it?",
     "Workarounds,\nparallel workflows",
     "Same workflow,\nconsistent use"),
    ("ACCOUNTABILITY",  "Is it really working?",
     "Different data,\ndifferent stories",
     "Same data,\nsame story"),
    ("ASSUMPTIONS",     "What are we assuming?",
     "Model trusted\nover reality",
     "Reality tested,\nassumptions exposed"),
]

dw = (W - 2*M - Inches(0.45)) / 4
dh = Inches(4.2)
d_gap = Inches(0.15)
d_top = Inches(1.55)

for i, (dim, q, red, green) in enumerate(dims):
    l = M + i * (dw + d_gap)
    rect(s8, l, d_top, dw, dh, LGRAY, MGRAY, 1)
    rect(s8, l, d_top, dw, Inches(0.75), BLACK)
    tb(s8, dim, l + Inches(0.12), d_top + Inches(0.07), dw - Inches(0.24), Inches(0.38),
       size=13, bold=True, color=WHITE)
    tb(s8, q, l + Inches(0.12), d_top + Inches(0.45), dw - Inches(0.24), Inches(0.35),
       size=11, italic=True, color=MGRAY)

    # Red light
    rect(s8, l + Inches(0.12), d_top + Inches(0.95), Inches(0.25), Inches(0.25),
         RGBColor(0xCC,0x00,0x00))
    tb(s8, "RED LIGHT", l + Inches(0.45), d_top + Inches(0.95), dw - Inches(0.57), Inches(0.3),
       size=10, bold=True, color=RED)
    tb(s8, red, l + Inches(0.12), d_top + Inches(1.3), dw - Inches(0.24), Inches(0.9),
       size=11, color=DGRAY)

    # Green light
    rect(s8, l + Inches(0.12), d_top + Inches(2.35), Inches(0.25), Inches(0.25),
         RGBColor(0x2E,0x7D,0x32))
    tb(s8, "GREEN LIGHT", l + Inches(0.45), d_top + Inches(2.35), dw - Inches(0.57), Inches(0.3),
       size=10, bold=True, color=RGBColor(0x2E,0x7D,0x32))
    tb(s8, green, l + Inches(0.12), d_top + Inches(2.7), dw - Inches(0.24), Inches(0.9),
       size=11, color=DGRAY)

# Bottom CTA bar
rect(s8, M, Inches(5.9), W - 2*M, Inches(0.7), BLACK)
tb(s8, "START WITH YOUR LOWEST SCORE — that's your highest-leverage leadership action.",
   M + Inches(0.2), Inches(6.0), W - 2*M - Inches(0.4), Inches(0.5),
   size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

notes(s8, "Walk each quadrant. This tool scores any AI initiative from 0–8. A score of 0-2 = serious execution risk. 3-5 = mixed, targeted intervention needed. 6-8 = healthy, maintain. The instruction: start with the lowest score, not the average.")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — If You See This, Do This
# ─────────────────────────────────────────────────────────────────────────────
s9 = prs.slides.add_slide(blank)
bg(s9, WHITE)

rect(s9, 0, 0, W, Inches(0.1), RED)

tb(s9, "IF YOU SEE THIS, DO THIS",
   M, Inches(0.25), W - 2*M, Inches(0.75),
   size=34, bold=True, color=BLACK)

tb(s9, "Concrete actions for each execution gap",
   M, Inches(0.95), W - 2*M, Inches(0.4),
   size=15, italic=True, color=DGRAY)

rows = [
    ("VISIBILITY GAP",      "Fragmented inputs, hidden dependencies",
     "• Map inputs end-to-end (who creates, changes, consumes)\n• Assign ownership for each dependency — no invisible handoffs"),
    ("ADOPTION GAP",        "Workarounds, parallel workflows",
     "• Standardize one way of working (remove optional paths)\n• Make AI outputs explainable — users understand why, not just what"),
    ("ACCOUNTABILITY GAP",  "Different data, different stories",
     "• Define one shared metric of success (visible to all stakeholders)\n• Review results in one forum — same data, same conversation"),
    ("ASSUMPTION GAP",      "Model trusted over reality",
     "• Test outputs in real conditions, not just system data\n• Create weekly feedback loop from frontline users — not ad hoc"),
]

row_h = Inches(1.2)
r_top = Inches(1.5)
col1_w = Inches(2.8)
col2_w = Inches(3.5)
col3_w = W - 2*M - col1_w - col2_w - Inches(0.3)

for i, (gap, symptom, action) in enumerate(rows):
    t = r_top + i * (row_h + Inches(0.1))
    bg_c = LGRAY if i % 2 == 0 else OFFWHITE
    rect(s9, M, t, W - 2*M, row_h, bg_c)
    rect(s9, M, t, Inches(0.06), row_h, RED)
    tb(s9, gap, M + Inches(0.15), t + Inches(0.1), col1_w, Inches(0.45),
       size=13, bold=True, color=BLACK)
    tb(s9, symptom, M + Inches(0.15), t + Inches(0.55), col1_w, Inches(0.55),
       size=11, italic=True, color=DGRAY)
    tb(s9, action, M + col1_w + Inches(0.25), t + Inches(0.15),
       col2_w + col3_w, row_h - Inches(0.3),
       size=12, color=BLACK)

# Footer
rect(s9, M, H - Inches(0.6), W - 2*M, Inches(0.04), RED)
tb(s9, "AI succeeds where these conversations happen. It fails where they don't.",
   M, H - Inches(0.55), W - 2*M, Inches(0.45),
   size=13, bold=True, italic=True, color=RED, align=PP_ALIGN.CENTER)

notes(s9, "This is the action slide. For each gap, the fix is a conversation — not a technology update. Walk through each row and ask: 'Which of these is present in your current AI initiatives?' The bottom tagline is the closing thesis of the entire session.")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Call to Action
# ─────────────────────────────────────────────────────────────────────────────
s10 = prs.slides.add_slide(blank)
bg(s10, BLACK)

rect(s10, 0, 0, W, Inches(0.1), RED)

rect(s10, M, Inches(0.5), Inches(3), Inches(0.55), RED)
tb(s10, "OPERATOR PERSPECTIVE", M + Inches(0.12), Inches(0.57), Inches(2.8), Inches(0.42),
   size=12, bold=True, color=WHITE)

tb(s10, "WHERE IS YOUR\nAI BREAKING?",
   M, Inches(1.3), Inches(5.5), Inches(2.2),
   size=44, bold=True, color=WHITE)

tb(s10, "That's your leadership opportunity.",
   M, Inches(3.5), Inches(5.5), Inches(0.7),
   size=22, italic=True, color=RED)

# Right: 4 bullets
rect(s10, Inches(7.0), Inches(0.5), Inches(5.8), Inches(5.6), RGBColor(0x1A, 0x1A, 0x1A))

multiline_tb(s10, [
    ("Diagnose against the four dimensions:", True, WHITE),
    ("", False, WHITE),
    ("•  Visibility — what are we not seeing?", False, WHITE),
    ("•  Adoption — are people actually using it?", False, WHITE),
    ("•  Accountability — does one shared story exist?", False, WHITE),
    ("•  Assumptions — is the model grounded in reality?", False, WHITE),
    ("", False, WHITE),
    ("Use this diagnostic to see it faster and fix it.", True, RED),
],
   Inches(7.2), Inches(0.75), Inches(5.4), Inches(5.1),
   size=16, color=WHITE)

# Bottom bar
rect(s10, 0, Inches(6.3), W, Inches(1.2), RED)
tb(s10,
   "AI promises unprecedented intelligence.\nYour organization's agile people will turn that intelligence into results.",
   M, Inches(6.4), W - 2*M, Inches(1.0),
   size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

notes(s10, "Close by returning to the opening question: 'The model was right. The outcome was wrong. Where is that happening in your organization right now?' End with the diagnostic as the next step — something they can do today.")


# ── save ──────────────────────────────────────────────────────────────────────
OUT = "AI-That-Delivers-Summary.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
