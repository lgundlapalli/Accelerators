"""
Generate native editable PPTX for 'Why AI Initiatives Fail' presentation.
Abbott BTS-DTS brand: Deep Navy (#000050), Lavender (#DDDDF8), White, Near-Black.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
NAVY = RGBColor(0x00, 0x00, 0x50)
LAVENDER = RGBColor(0xDD, 0xDD, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
MED_BLUE = RGBColor(0x33, 0x33, 0xAA)
LIGHT_GRAY = RGBColor(0x88, 0x88, 0x88)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_footer(slide, page_num):
    """Add confidentiality footer and page number."""
    tf = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(8), Inches(0.4)).text_frame
    p = tf.paragraphs[0]
    p.text = "Proprietary and confidential — do not distribute"
    p.font.size = Pt(9)
    p.font.color.rgb = LIGHT_GRAY
    p.font.italic = True

    tf2 = slide.shapes.add_textbox(Inches(12), Inches(6.9), Inches(1), Inches(0.4)).text_frame
    p2 = tf2.paragraphs[0]
    p2.text = str(page_num)
    p2.font.size = Pt(9)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.RIGHT


def navy_bg(slide):
    """Set slide background to navy."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY


def add_text(slide, left, top, width, height, text, font_size=18, bold=False,
             color=NEAR_BLACK, alignment=PP_ALIGN.LEFT, italic=False, font_name='Helvetica Neue'):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.font.italic = italic
    p.font.name = font_name
    return tf


def add_multiline(slide, left, top, width, height, lines, font_size=14,
                   color=NEAR_BLACK, bold_first=False, line_spacing=1.2):
    """Add multi-line text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
        p.font.name = 'Helvetica Neue'
        if bold_first and i == 0:
            p.font.bold = True
    return tf


def add_table(slide, left, top, width, height, rows, cols, data,
              col_widths=None, header_color=LAVENDER, font_size=11):
    """Add a formatted table."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = 'Helvetica Neue'
                paragraph.font.color.rgb = NEAR_BLACK
                paragraph.space_after = Pt(2)

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = NAVY
            else:
                cell.fill.background()

    return table


# ============================================================
# SLIDE 1: Title (Navy Background)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
navy_bg(slide)

add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.2),
         "Why AI Initiatives Fail", font_size=44, bold=False, color=WHITE)
add_text(slide, Inches(0.8), Inches(2.7), Inches(11), Inches(1),
         "Cultural Debt Is Your Biggest AI Risk", font_size=40, bold=False, color=WHITE)
add_text(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(0.6),
         "The hidden organizational drag on AI ROI — and how to unwind it | BTS-DTS Architecture",
         font_size=16, italic=True, color=RGBColor(0xBB, 0xBB, 0xFF))

# ============================================================
# SLIDE 2: Executive Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Executive Summary", font_size=32, bold=False, color=NAVY)

add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8),
         "AI initiatives fail at a 95% rate — not because of technology, but because organizations deploy 2026 tools inside 2019 cultures. Cultural debt is the silent compounding force behind every stalled pilot and every missed ROI target.",
         font_size=13, color=NEAR_BLACK)

data = [
    ["The Problem", "The Root Cause", "The Fix"],
    ["1 in 4", "Cultural Debt", "Behavioral Agility"],
    ["Only 25% of AI initiatives deliver expected ROI. Technology works; adoption doesn't.",
     "Unresolved habits, unexamined processes, and unspoken assumptions silently compound — encoding inefficiency into every AI workflow.",
     "Organizations that align all six sources of behavioral influence are 10x more likely to achieve rapid change."],
    ["$Billions Wasted", "Bureaucratic Antibodies", "Treat It Like Technical Debt"],
    ["90% of organizations use AI, yet few report scaled enterprise impact. The gap is widening.",
     "Middle managers whose authority depends on current workflows predictably sabotage pilots — not maliciously, but rationally.",
     "Audit cultural debt. Redesign incentives. Rebuild workflows from the outcome backward."],
]
add_table(slide, Inches(0.6), Inches(1.9), Inches(12), Inches(4.2), 5, 3, data,
          col_widths=[Inches(4), Inches(4), Inches(4)], font_size=11)

add_text(slide, Inches(0.6), Inches(6.2), Inches(12), Inches(0.4),
         "Sources: IBM Global AI Survey (2025); McKinsey State of AI (2025); HBR (2025); Shelly Palmer (2026)",
         font_size=9, italic=True, color=LIGHT_GRAY)
add_footer(slide, 2)

# ============================================================
# SLIDE 3: Agenda
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Agenda", font_size=32, bold=False, color=NAVY)

agenda_items = [
    ("01", "The ROI Gap — What the Data Shows"),
    ("02", "Cultural Debt — The Hidden Compounding Force"),
    ("03", "Bureaucratic Antibodies — How Pilots Get Killed"),
    ("04", "Six Sources of Influence — The Behavioral Lever"),
    ("05", "Unwinding Cultural Debt — The Path Forward"),
]

for i, (num, text) in enumerate(agenda_items):
    y = Inches(1.6) + Inches(i * 0.9)
    add_text(slide, Inches(1.2), y, Inches(1.2), Inches(0.6),
             num, font_size=28, bold=False, color=NAVY)
    add_text(slide, Inches(2.8), y + Inches(0.05), Inches(8), Inches(0.6),
             text, font_size=18, color=NEAR_BLACK)
    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), y + Inches(0.7), Inches(10), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    line.line.fill.background()

add_footer(slide, 3)

# ============================================================
# SLIDE 4: Section Divider — The ROI Gap
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
navy_bg(slide)
add_text(slide, Inches(0.8), Inches(2.5), Inches(4), Inches(0.8),
         "01", font_size=48, bold=False, color=WHITE)
add_text(slide, Inches(0.8), Inches(3.5), Inches(10), Inches(1),
         "The ROI Gap", font_size=40, bold=False, color=WHITE)

# ============================================================
# SLIDE 5: Big Number Impact Grid
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Only 1 in 4 AI Initiatives Deliver Expected ROI", font_size=28, bold=False, color=NAVY)
add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.5),
         "The technology is not the bottleneck. Organizational readiness is.",
         font_size=14, color=NEAR_BLACK)

metrics = [
    ("1 in 4", "Deliver Expected ROI", "IBM: only 25% of AI initiatives meet ROI targets"),
    ("95%", "Fail to Deliver Value", "HBR: failures from \"technosolutionism\" — assuming better tools solve org problems"),
    ("90%", "Use AI Somewhere", "McKinsey: deployed everywhere, scaled nowhere"),
    ("82%", "Use GenAI Weekly", "Wharton: leaders use it regularly, orgs haven't changed"),
]

for i, (number, label, desc) in enumerate(metrics):
    x = Inches(0.6) + Inches(i * 3.1)
    # Number
    add_text(slide, x, Inches(1.8), Inches(2.8), Inches(1),
             number, font_size=44, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
    # Label
    add_text(slide, x, Inches(2.9), Inches(2.8), Inches(0.5),
             label, font_size=14, bold=True, color=MED_BLUE, alignment=PP_ALIGN.CENTER)
    # Description
    add_text(slide, x, Inches(3.5), Inches(2.8), Inches(1.5),
             desc, font_size=11, color=NEAR_BLACK, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0.6), Inches(5.4), Inches(12), Inches(0.5),
         "The pattern: AI tools are advancing rapidly. Human systems are not.",
         font_size=16, bold=True, color=NAVY)

add_text(slide, Inches(0.6), Inches(6.2), Inches(12), Inches(0.4),
         "Sources: IBM/Fortune (2025); HBR, De Cremer et al. (2025); McKinsey Global Survey (2025); Wharton (2025)",
         font_size=9, italic=True, color=LIGHT_GRAY)
add_footer(slide, 4)

# ============================================================
# SLIDE 6: Section Divider — Cultural Debt
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
navy_bg(slide)
add_text(slide, Inches(0.8), Inches(2.5), Inches(4), Inches(0.8),
         "02", font_size=48, bold=False, color=WHITE)
add_text(slide, Inches(0.8), Inches(3.5), Inches(10), Inches(1),
         "Cultural Debt", font_size=40, bold=False, color=WHITE)

# ============================================================
# SLIDE 7: Cultural Debt Defined
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Cultural Debt Compounds Silently Until It Becomes Your Biggest AI Constraint",
         font_size=24, bold=False, color=NAVY)

data = [
    ["", "Technical Debt", "Cultural Debt"],
    ["Definition", "Outdated infrastructure, legacy code, integration gaps",
     "Unresolved habits, unexamined processes, unspoken assumptions"],
    ["Visibility", "Visible in architecture reviews and incident reports",
     "Invisible — masquerades as \"how we do things\""],
    ["Who Fixes It", "Engineers can measure and fix it",
     "Leadership must surface and confront it"],
    ["Solution", "Modernization and refactoring",
     "Incentive redesign and behavioral change"],
    ["Impact on AI", "Slows deployment",
     "Prevents adoption"],
]
add_table(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(3.2), 6, 3, data,
          col_widths=[Inches(2.5), Inches(4.75), Inches(4.75)], font_size=12)

# Quote box
quote_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.7), Inches(12), Inches(1))
quote_box.fill.solid()
quote_box.fill.fore_color.rgb = LAVENDER
quote_box.line.fill.background()

add_text(slide, Inches(1), Inches(4.85), Inches(11), Inches(0.7),
         "\"If you are asking a 2026 technology to operate inside a 2019 culture, the culture will win every time.\" — Shelly Palmer",
         font_size=16, bold=True, color=NAVY, italic=True)

add_text(slide, Inches(0.6), Inches(5.9), Inches(12), Inches(0.5),
         "AI agents are ruthlessly literal. They encode and deepen cultural debt: the unnecessary sign-offs, the hedging language from a forgotten complaint, the approval chain no one questions.",
         font_size=12, color=NEAR_BLACK)

add_footer(slide, 5)

# ============================================================
# SLIDE 8: Section Divider — Bureaucratic Antibodies
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
navy_bg(slide)
add_text(slide, Inches(0.8), Inches(2.5), Inches(4), Inches(0.8),
         "03", font_size=48, bold=False, color=WHITE)
add_text(slide, Inches(0.8), Inches(3.5), Inches(10), Inches(1),
         "Bureaucratic Antibodies", font_size=40, bold=False, color=WHITE)

# ============================================================
# SLIDE 9: Antibodies Kill Pilots
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Bureaucratic Antibodies Kill AI Pilots in 3 Predictable Steps",
         font_size=24, bold=False, color=NAVY)

add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.5),
         "Process owners whose authority depends on current workflows are not saboteurs — their incentives are entirely rational.",
         font_size=13, color=NEAR_BLACK)

steps = [
    ("Step 1: Narrow the Scope", "The pilot gets scoped so narrowly it cannot demonstrate value."),
    ("Step 2: Expand the Criteria", "AI is measured against standards no human employee has ever met."),
    ("Step 3: Declare \"Promising but Premature\"", "Everyone returns to the process they already know. Cultural debt compounds."),
]

for i, (title, desc) in enumerate(steps):
    x = Inches(0.6) + Inches(i * 4.1)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.8), Inches(3.8), Inches(1.8))
    box.fill.solid()
    box.fill.fore_color.rgb = LAVENDER
    box.line.fill.background()
    add_text(slide, x + Inches(0.2), Inches(1.95), Inches(3.4), Inches(0.5),
             title, font_size=14, bold=True, color=NAVY)
    add_text(slide, x + Inches(0.2), Inches(2.5), Inches(3.4), Inches(0.9),
             desc, font_size=12, color=NEAR_BLACK)

# Sinclair quote
quote_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.0), Inches(12), Inches(0.8))
quote_box.fill.solid()
quote_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xFF)
quote_box.line.fill.background()

add_text(slide, Inches(1), Inches(4.1), Inches(11), Inches(0.6),
         "\"It is difficult to get a man to understand something when his job depends on him not understanding it.\" — Upton Sinclair",
         font_size=14, italic=True, color=NAVY)

add_text(slide, Inches(0.6), Inches(5.1), Inches(12), Inches(1),
         "The Electrification Parallel: Electric power was demonstrated in 1880. Factory productivity gains didn't materialize until the 1920s. The lag was not technological — factories had to be physically rebuilt, and the foremen retrained. We are in the same reorganization gap now, except the reorganization is cultural.",
         font_size=12, color=NEAR_BLACK)

add_text(slide, Inches(0.6), Inches(6.2), Inches(12), Inches(0.4),
         "Source: Shelly Palmer, \"Unwinding Cultural Debt\" (2026)",
         font_size=9, italic=True, color=LIGHT_GRAY)
add_footer(slide, 6)

# ============================================================
# SLIDE 10: Section Divider — Six Sources
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
navy_bg(slide)
add_text(slide, Inches(0.8), Inches(2.5), Inches(4), Inches(0.8),
         "04", font_size=48, bold=False, color=WHITE)
add_text(slide, Inches(0.8), Inches(3.5), Inches(10), Inches(1),
         "Six Sources of Influence", font_size=40, bold=False, color=WHITE)

# ============================================================
# SLIDE 11: Six Sources Table
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Behavioral Agility Requires Aligning All Six Sources of Influence",
         font_size=24, bold=False, color=NAVY)

add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.5),
         "Fastest adopters are 12x quicker at turning ideas into habits. The six-source approach increases odds of rapid change 10x.",
         font_size=13, color=NEAR_BLACK)

data = [
    ["", "Motivation", "Ability"],
    ["Personal",
     "Reframe AI from \"threat\" to \"strategic advantage\" through storytelling and direct experience",
     "Hands-on workshops: practice prompting, validating, integrating AI into real workflows"],
    ["Social",
     "Respected leaders visibly model AI use — shift from \"optional innovation\" to \"how high performers operate\"",
     "Cross-functional AI councils: reduce fear, increase trust, accelerate responsible adoption"],
    ["Structural",
     "Tie AI adoption to promotion criteria — reward adoption, not just experimentation",
     "Embed AI into existing workflows — make adoption automatic; require HITL eval as job function"],
]
add_table(slide, Inches(0.6), Inches(1.7), Inches(12), Inches(3.5), 4, 3, data,
          col_widths=[Inches(2), Inches(5), Inches(5)], font_size=12)

# Bottom callout
callout = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.5), Inches(12), Inches(0.7))
callout.fill.solid()
callout.fill.fore_color.rgb = LAVENDER
callout.line.fill.background()

add_text(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.5),
         "If even one source is misaligned, adoption stalls. If several are misaligned, it fails entirely.",
         font_size=15, bold=True, color=NAVY)

add_text(slide, Inches(0.6), Inches(6.4), Inches(12), Inches(0.3),
         "Sources: Crucial Learning (2026); Grenny, Maxfield, Shimberg — MIT Sloan Management Review (2008)",
         font_size=9, italic=True, color=LIGHT_GRAY)
add_footer(slide, 7)

# ============================================================
# SLIDE 12: Section Divider — Path Forward
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
navy_bg(slide)
add_text(slide, Inches(0.8), Inches(2.5), Inches(4), Inches(0.8),
         "05", font_size=48, bold=False, color=WHITE)
add_text(slide, Inches(0.8), Inches(3.5), Inches(10), Inches(1),
         "The Path Forward", font_size=40, bold=False, color=WHITE)

# ============================================================
# SLIDE 13: Phased Roadmap
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Unwind Cultural Debt With the Same Rigor as Technical Debt",
         font_size=24, bold=False, color=NAVY)

add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
         "ROI on AI always accurately reflects what an organization's cultural debt allows.",
         font_size=14, bold=True, color=MED_BLUE)

phases = [
    ("Phase 1 — Audit", "Now — This Quarter",
     "Surface hidden cultural debt",
     "Map every process AI will touch. Identify \"antibodies\" — roles whose authority depends on current workflows. Ask: \"Would this step exist if we built from scratch?\"",
     "Decision",
     "Approve cultural debt audit as mandatory companion to every AI-First initiative"),
    ("Phase 2 — Redesign", "Months 2–3",
     "Realign incentives and workflows",
     "Redesign incentive structures: reward adoption, not experimentation. Pay people to learn evals and HITL review. Rebuild workflows from outcome backward.",
     "Success Signal",
     "Adoption rate up, antibody resistance down, delivery speed improvement"),
    ("Phase 3 — Embed", "Month 4 Onwards",
     "Make behavioral agility permanent",
     "Integrate Six Sources into every AI initiative launch. Build cultural debt audits into POD operating cadence. Measure behavioral change alongside delivery.",
     "Long-term Value",
     "AI ROI reflects organizational capacity, not just technical capability"),
]

for i, (title, timing, headline, details, signal_label, signal_text) in enumerate(phases):
    x = Inches(0.6) + Inches(i * 4.1)
    # Phase box
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.6), Inches(3.8), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = LAVENDER
    box.line.fill.background()

    add_text(slide, x + Inches(0.2), Inches(1.7), Inches(3.4), Inches(0.4),
             title, font_size=16, bold=True, color=NAVY)
    add_text(slide, x + Inches(0.2), Inches(2.1), Inches(3.4), Inches(0.3),
             timing, font_size=11, italic=True, color=MED_BLUE)
    add_text(slide, x + Inches(0.2), Inches(2.5), Inches(3.4), Inches(0.4),
             headline, font_size=13, bold=True, color=NEAR_BLACK)
    add_text(slide, x + Inches(0.2), Inches(3.0), Inches(3.4), Inches(1.5),
             details, font_size=10, color=NEAR_BLACK)
    add_text(slide, x + Inches(0.2), Inches(4.7), Inches(3.4), Inches(0.3),
             signal_label, font_size=11, bold=True, color=NAVY)
    add_text(slide, x + Inches(0.2), Inches(5.1), Inches(3.4), Inches(0.8),
             signal_text, font_size=10, color=NEAR_BLACK)

add_footer(slide, 8)

# ============================================================
# SLIDE 14: Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Three Actions This Quarter", font_size=28, bold=False, color=NAVY)

actions = [
    ("1. Add a Cultural Debt Audit to every AI-First POD kickoff",
     "Surface unexamined processes before encoding them into agents.",
     "Owner: Architecture + POD Leads  |  By: End of quarter"),
    ("2. Redesign one incentive structure",
     "Tie responsible AI adoption to performance reviews and promotion criteria.",
     "Owner: HR + Transformation Steering Committee  |  By: Next performance cycle"),
    ("3. Run a \"Future-Back\" workflow design session",
     "Force one POD to design assuming AI is primary producer, humans evaluate quality.",
     "Owner: AI Workflow Orchestrators  |  By: Within 30 days"),
]

for i, (title, desc, owner) in enumerate(actions):
    y = Inches(1.3) + Inches(i * 1.5)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(12), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = LAVENDER if i % 2 == 0 else RGBColor(0xF0, 0xF0, 0xFF)
    box.line.fill.background()

    add_text(slide, Inches(1), y + Inches(0.1), Inches(11), Inches(0.4),
             title, font_size=15, bold=True, color=NAVY)
    add_text(slide, Inches(1), y + Inches(0.5), Inches(11), Inches(0.35),
             desc, font_size=12, color=NEAR_BLACK)
    add_text(slide, Inches(1), y + Inches(0.9), Inches(11), Inches(0.3),
             owner, font_size=10, bold=True, color=MED_BLUE)

# Closing quote
add_text(slide, Inches(0.6), Inches(5.9), Inches(12), Inches(0.6),
         "The question is not whether you've invested in AI. The question is whether you've invested in the human system required to make it work.",
         font_size=14, bold=True, color=NAVY, italic=True)

add_footer(slide, 9)

# ============================================================
# SLIDE 15: Appendix Divider
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
navy_bg(slide)
add_text(slide, Inches(0.8), Inches(3), Inches(10), Inches(1.2),
         "Appendix", font_size=44, bold=False, color=WHITE)

# ============================================================
# SLIDE 16: Anticipated Questions
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Anticipated Questions", font_size=28, bold=False, color=NAVY)

questions = [
    ("\"Isn't this just change management?\"",
     "Traditional change management covers 2 of 6 influence sources. That's why it fails. Cultural debt requires all six."),
    ("\"How do we measure cultural debt?\"",
     "Three proxies: (1) Time from pilot to production. (2) Process steps no one can explain. (3) Gap between AI deployed vs. AI used weekly."),
    ("\"Won't this slow down AI-First delivery?\"",
     "The opposite. Cultural debt IS the bottleneck. Every \"promising but premature\" verdict costs more than an audit."),
    ("\"Are middle managers the problem?\"",
     "No. Their incentives are. Redesign the incentives and the antibody response disappears."),
    ("\"How does this connect to our AI-First Operating Model?\"",
     "The operating model defines WHAT. Cultural debt unwinding defines WHETHER it gets adopted."),
]

for i, (q, a) in enumerate(questions):
    y = Inches(1.1) + Inches(i * 1.1)
    add_text(slide, Inches(0.6), y, Inches(5), Inches(0.4),
             q, font_size=12, bold=True, color=NAVY)
    add_text(slide, Inches(5.8), y, Inches(7), Inches(0.9),
             a, font_size=11, color=NEAR_BLACK)

add_footer(slide, 10)

# ============================================================
# SLIDE 17: Key Statistics
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Key Statistics Reference", font_size=28, bold=False, color=NAVY)

data = [
    ["Statistic", "Source", "Year"],
    ["Only 1 in 4 AI initiatives deliver expected ROI", "IBM Global AI Survey / Fortune", "2025"],
    ["95% of AI initiatives fail to deliver value", "HBR — De Cremer, Schweitzer et al.", "2025"],
    ["90% of orgs use AI in at least one function", "McKinsey State of AI Global Survey", "2025"],
    ["82% of leaders use GenAI weekly", "Wharton AI Adoption Report", "2025"],
    ["12x faster adoption in agile organizations", "Crucial Learning (1,700+ professionals)", "2026"],
    ["10x higher odds with six-source approach", "MIT Sloan Management Review", "2008"],
    ["Electrification: 40-year reorganization gap", "Economic historians", "Historical"],
]
add_table(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(4), 8, 3, data,
          col_widths=[Inches(5.5), Inches(4.5), Inches(2)], font_size=12)

add_footer(slide, 11)

# ============================================================
# SAVE
# ============================================================
output_path = "/Users/GUNDLLX/learn-claude/presentations/why-ai-initiatives-fail-editable.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
