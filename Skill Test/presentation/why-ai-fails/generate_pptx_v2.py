"""
'Why AI Initiatives Fail' — V2: Anchored on Six Sources of Influence.
Cultural debt is diagnosed AND solved through the Six Sources framework.
Abbott BTS-DTS brand.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
NAVY = RGBColor(0x00, 0x00, 0x50)
LAVENDER = RGBColor(0xDD, 0xDD, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
MED_BLUE = RGBColor(0x33, 0x33, 0xAA)
LIGHT_GRAY = RGBColor(0x88, 0x88, 0x88)
SOFT_LAVENDER = RGBColor(0xF0, 0xF0, 0xFF)
RED_ACCENT = RGBColor(0xCC, 0x33, 0x33)
GREEN_ACCENT = RGBColor(0x22, 0x88, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_footer(slide, page_num):
    tf = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(8), Inches(0.4)).text_frame
    p = tf.paragraphs[0]
    p.text = "Proprietary and confidential — do not distribute"
    p.font.size = Pt(9)
    p.font.color.rgb = LIGHT_GRAY
    p.font.italic = True
    tf2 = slide.shapes.add_textbox(Inches(12), Inches(6.9), Inches(1), Inches(0.4)).text_frame
    p2 = tf2.paragraphs[0]
    p2.text = str(page_num)
    p2.font.size = Pt(9)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.RIGHT


def navy_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY


def add_text(slide, left, top, width, height, text, font_size=18, bold=False,
             color=NEAR_BLACK, alignment=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.font.italic = italic
    p.font.name = 'Helvetica Neue'
    return tf


def add_rich_text(slide, left, top, width, height, segments):
    """Add text with mixed formatting. segments = [(text, size, bold, color, italic), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    for i, (text, size, bold, color, italic) in enumerate(segments):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.italic = italic
        run.font.name = 'Helvetica Neue'
    return tf


def add_box(slide, left, top, width, height, fill_color):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.fill.background()
    return box


def add_table(slide, left, top, width, height, rows, cols, data,
              col_widths=None, font_size=11, header_color=LAVENDER):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = 'Helvetica Neue'
                paragraph.font.color.rgb = NEAR_BLACK
                paragraph.space_after = Pt(2)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = NAVY
            else:
                cell.fill.background()
    return table


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
navy_bg(slide)
add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(1.2),
         "Why AI Initiatives Fail", font_size=44, color=WHITE)
add_text(slide, Inches(0.8), Inches(2.4), Inches(11), Inches(1),
         "And How Cultural Debt Silently Kills Your AI ROI", font_size=36, color=WHITE)
add_text(slide, Inches(0.8), Inches(4.0), Inches(10), Inches(0.6),
         "A Six Sources of Influence perspective on what's really blocking adoption",
         font_size=16, italic=True, color=RGBColor(0xBB, 0xBB, 0xFF))
add_text(slide, Inches(0.8), Inches(5.5), Inches(10), Inches(0.4),
         "BTS-DTS Architecture", font_size=14, color=RGBColor(0x99, 0x99, 0xCC))

add_notes(slide, """SPEAKING POINTS — Title Slide

OPENING (15 seconds):
- "We've made significant AI investments. Today I want to talk about why most of them — across the industry, not just here — aren't delivering the returns we expected. And more importantly, what we can do about it."
- Don't read the slide. Set the room with a direct statement.

CONTEXT:
- This presentation draws on two key sources: a Crucial Learning research report on AI initiative failure (surveying 1,700+ professionals), and Shelly Palmer's analysis of cultural debt in enterprise AI adoption.
- The central argument: AI failure is not a technology problem. It's a human system problem. And there's a well-researched framework — Six Sources of Influence — that diagnoses exactly where adoption breaks and how to fix it.

TRANSITION:
- "Let me start with the numbers, because they're sobering."
""")

# ============================================================
# SLIDE 2: Executive Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "Executive Summary", font_size=32, color=NAVY)

add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.7),
         "95% of AI initiatives fail — not from bad technology, but from cultural debt: the unresolved habits, broken incentives, and silent resistance that block adoption. The Six Sources of Influence framework diagnoses exactly where adoption breaks and prescribes how to fix it.",
         font_size=13, color=NEAR_BLACK)

data = [
    ["The Problem", "The Root Cause", "The Framework"],
    ["1 in 4 deliver expected ROI", "Cultural Debt compounds silently", "Six Sources of Influence"],
    ["Technology works. Behavior doesn't change. 90% of orgs deploy AI; few see scaled impact.",
     "Unresolved habits, misaligned incentives, and \"bureaucratic antibodies\" encode dysfunction into AI workflows.",
     "Behavior is shaped by 6 forces: personal, social, and structural — each with motivation and ability. Miss one, adoption stalls."],
    ["$Billions Wasted", "Invisible Until It's Too Late", "10x Higher Odds of Success"],
    ["Every stalled pilot and \"promising but premature\" verdict is cultural debt collecting interest.",
     "It masquerades as \"how we do things\" — nobody questions it until AI makes it visible.",
     "Leaders who engage all 6 sources are 10x more likely to drive rapid behavioral change."],
]
add_table(slide, Inches(0.6), Inches(1.7), Inches(12), Inches(4.2), 5, 3, data,
          col_widths=[Inches(4), Inches(4), Inches(4)], font_size=11)

add_text(slide, Inches(0.6), Inches(6.1), Inches(12), Inches(0.4),
         "Sources: IBM/Fortune (2025); McKinsey (2025); HBR (2025); Crucial Learning (2026); Shelly Palmer (2026)",
         font_size=9, italic=True, color=LIGHT_GRAY)
add_footer(slide, 2)

add_notes(slide, """SPEAKING POINTS — Executive Summary

PURPOSE: This is the "read only this slide" slide. If someone walks in late or skips ahead, they get the full story here.

KEY TALKING POINTS:
- "Three columns tell the whole story: the problem is clear — 1 in 4 AI initiatives deliver ROI. The root cause is not technology — it's cultural debt, the accumulated habits and assumptions that block adoption. And the fix is a research-backed framework called Six Sources of Influence."
- Emphasize: "This isn't opinion. These are findings from IBM, McKinsey, Harvard Business Review, and a 1,700-person research study."
- The word "cultural debt" is intentional — it parallels "technical debt" which this audience already understands. Technical debt slows deployment. Cultural debt prevents adoption entirely.

ANTICIPATED PUSHBACK:
- "We already do change management" → Traditional change management covers 2 of 6 influence sources. That's why it has a 70% failure rate. We'll get to that.

TRANSITION:
- "Let me walk you through the five sections."
""")

# ============================================================
# SLIDE 3: Agenda
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Agenda", font_size=32, color=NAVY)

agenda = [
    ("01", "The Numbers — Why 95% of AI Initiatives Fail"),
    ("02", "Cultural Debt — The Hidden Force Behind the Failures"),
    ("03", "The Six Sources Framework — Where Adoption Breaks"),
    ("04", "Diagnosing Your Cultural Debt Through Six Sources"),
    ("05", "Unwinding It — Source by Source"),
]
for i, (num, text) in enumerate(agenda):
    y = Inches(1.5) + Inches(i * 0.95)
    add_text(slide, Inches(1.2), y, Inches(1.2), Inches(0.6), num, font_size=28, color=NAVY)
    add_text(slide, Inches(2.8), y + Inches(0.05), Inches(8), Inches(0.6), text, font_size=18, color=NEAR_BLACK)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), y + Inches(0.7), Inches(10), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    line.line.fill.background()
add_footer(slide, 3)

add_notes(slide, """SPEAKING POINTS — Agenda

- Don't spend more than 15 seconds here. Simply say: "Five sections — we start with the data, then define the problem, introduce the framework, use it to diagnose where we're stuck, and end with three concrete actions."
- This is a 10-minute presentation. Keep momentum. Don't narrate the agenda — move to the numbers.
""")

# ============================================================
# SLIDE 4: Divider — The Numbers
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
navy_bg(slide)
add_text(slide, Inches(0.8), Inches(2.5), Inches(4), Inches(0.8), "01", font_size=48, color=WHITE)
add_text(slide, Inches(0.8), Inches(3.5), Inches(10), Inches(1),
         "The Numbers", font_size=40, color=WHITE)

add_notes(slide, """SECTION DIVIDER — Pause briefly. Let the room reset before the data slide.""")

# ============================================================
# SLIDE 5: Big Numbers
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Only 1 in 4 AI Initiatives Deliver Expected ROI", font_size=28, color=NAVY)
add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.5),
         "The technology is not the bottleneck. Organizational readiness is.",
         font_size=14, color=NEAR_BLACK)

metrics = [
    ("1 in 4", "Deliver Expected ROI", "IBM: despite massive investment, 75%\nmiss their ROI targets"),
    ("95%", "Fail to Deliver Value", "HBR: leaders default to\n\"technosolutionism\""),
    ("90%", "Use AI Somewhere", "McKinsey: deployed everywhere,\nscaled nowhere"),
    ("82%", "Use GenAI Weekly", "Wharton: leaders use it,\nbut orgs haven't changed"),
]
for i, (number, label, desc) in enumerate(metrics):
    x = Inches(0.6) + Inches(i * 3.1)
    add_text(slide, x, Inches(1.8), Inches(2.8), Inches(1),
             number, font_size=44, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(2.9), Inches(2.8), Inches(0.5),
             label, font_size=14, bold=True, color=MED_BLUE, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.5), Inches(2.8), Inches(1.2),
             desc, font_size=11, color=NEAR_BLACK, alignment=PP_ALIGN.CENTER)

add_box(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(0.7), LAVENDER)
add_text(slide, Inches(1), Inches(5.1), Inches(11), Inches(0.5),
         "AI tools are advancing rapidly. Human systems are not. The gap is where value dies.",
         font_size=16, bold=True, color=NAVY)

add_text(slide, Inches(0.6), Inches(6.2), Inches(12), Inches(0.3),
         "Sources: IBM/Fortune (2025); HBR, De Cremer et al. (2025); McKinsey (2025); Wharton (2025)",
         font_size=9, italic=True, color=LIGHT_GRAY)
add_footer(slide, 4)

add_notes(slide, """SPEAKING POINTS — The Numbers

DELIVERY: Walk left to right through each number. Let each one land.

1 in 4 (IBM):
- "IBM surveyed global leaders in 2025 and found only about 25 percent of AI initiatives delivered expected ROI. Three quarters missed. Not because the tools didn't work — because the organizations didn't change."

95% (HBR):
- "Harvard Business Review went further — up to 95 percent fail to deliver intended value. The authors coined the term 'technosolutionism': the assumption that better technology alone solves organizational problems. It doesn't."

90% (McKinsey):
- "McKinsey says 90 percent of organizations use AI somewhere. So it's not an access problem. It's a scaling problem. Deployed everywhere, scaled nowhere."

82% (Wharton):
- "Wharton found 82 percent of enterprise leaders personally use GenAI weekly. Leaders get it. But their organizations haven't followed."

THE PUNCHLINE:
- Point to the callout bar: "AI tools are advancing rapidly. Human systems are not. That gap is where the value dies."
- Pause. Let the room absorb.

TRANSITION:
- "So if the technology works, what's actually blocking us? The answer is cultural debt."
""")

# ============================================================
# SLIDE 6: Divider — Cultural Debt
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
navy_bg(slide)
add_text(slide, Inches(0.8), Inches(2.5), Inches(4), Inches(0.8), "02", font_size=48, color=WHITE)
add_text(slide, Inches(0.8), Inches(3.5), Inches(10), Inches(1),
         "Cultural Debt", font_size=40, color=WHITE)

add_notes(slide, """SECTION DIVIDER — "This is the concept that explains why the numbers look the way they do." """)

# ============================================================
# SLIDE 7: Cultural Debt vs Technical Debt
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Technical Debt Is an Engineering Problem. Cultural Debt Is a Leadership Problem.",
         font_size=24, color=NAVY)

# Two side-by-side panels
add_box(slide, Inches(0.6), Inches(1.2), Inches(5.8), Inches(3.5), SOFT_LAVENDER)
add_text(slide, Inches(0.9), Inches(1.3), Inches(5.2), Inches(0.4),
         "Technical Debt", font_size=20, bold=True, color=NAVY)

tech_items = [
    "Outdated infrastructure and legacy code",
    "Visible in architecture reviews",
    "Engineers can measure and fix it",
    "Solved by modernization",
    "Slows deployment",
]
for i, item in enumerate(tech_items):
    add_text(slide, Inches(0.9), Inches(1.85) + Inches(i * 0.45), Inches(5.2), Inches(0.4),
             f"  {item}", font_size=12, color=NEAR_BLACK)

add_box(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(3.5), LAVENDER)
add_text(slide, Inches(7.1), Inches(1.3), Inches(5.2), Inches(0.4),
         "Cultural Debt", font_size=20, bold=True, color=NAVY)

culture_items = [
    "Unresolved habits and unspoken assumptions",
    "Invisible — masquerades as \"how we do things\"",
    "Leadership must surface and confront it",
    "Solved by incentive redesign + behavioral change",
    "Prevents adoption entirely",
]
for i, item in enumerate(culture_items):
    add_text(slide, Inches(7.1), Inches(1.85) + Inches(i * 0.45), Inches(5.2), Inches(0.4),
             f"  {item}", font_size=12, color=NEAR_BLACK)

# Quote
add_box(slide, Inches(0.6), Inches(4.9), Inches(12), Inches(0.8), NAVY)
add_text(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.6),
         "\"If you are asking a 2026 technology to operate inside a 2019 culture, the culture will win every time.\" — Shelly Palmer",
         font_size=15, bold=True, color=WHITE, italic=True)

# Antibodies
add_text(slide, Inches(0.6), Inches(5.9), Inches(12), Inches(0.8),
         "\"Bureaucratic Antibodies\" — the predictable pattern: (1) Narrow the pilot scope so it can't prove value → (2) Expand eval criteria beyond what any human meets → (3) Declare \"promising but premature\" → everyone returns to what they know.",
         font_size=12, color=NEAR_BLACK)

add_footer(slide, 5)

add_notes(slide, """SPEAKING POINTS — Cultural Debt vs Technical Debt

FRAMING:
- "We all understand technical debt — it's in our vocabulary. Architecture reviews surface it. Engineers measure and fix it. But there's a parallel concept that's far more dangerous because it's invisible: cultural debt."

WALK THE COMPARISON:
- Point to the left panel: "Technical debt slows deployment. We know how to deal with it."
- Point to the right panel: "Cultural debt prevents adoption entirely. And nobody owns it."
- Key difference to emphasize: "Technical debt is visible in incident reports. Cultural debt masquerades as 'how we do things.' It's the weekly status meeting that exists because a VP got blindsided in 2019. The approval chain through a department that no longer does the work. The hiring rubric optimizing for credentials that mattered when we were solving a different problem."

THE QUOTE (Shelly Palmer):
- Read it aloud: "If you are asking a 2026 technology to operate inside a 2019 culture, the culture will win every time."
- Let it breathe. This is the thesis of the whole presentation.

ANTIBODIES:
- "Palmer calls the resistance forces 'bureaucratic antibodies.' These are not saboteurs — they are rational actors. If your job is to review every outbound proposal, and AI can draft one in four minutes, you have a choice: embrace the tool that eliminates your review function, or find seventeen reasons it needs more testing. Most people choose the testing."
- Describe the 3-step pattern: narrow scope → expand criteria → declare premature. "I've watched this pattern at organizations of every size."

TRANSITION:
- "So the question becomes: if the problem is behavioral, not technical, is there a framework for diagnosing and fixing it? There is. It's called the Six Sources of Influence."
""")

# ============================================================
# SLIDE 8: Divider — The Framework
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
navy_bg(slide)
add_text(slide, Inches(0.8), Inches(2.5), Inches(4), Inches(0.8), "03", font_size=48, color=WHITE)
add_text(slide, Inches(0.8), Inches(3.5), Inches(10), Inches(1),
         "The Six Sources Framework", font_size=40, color=WHITE)

add_notes(slide, """SECTION DIVIDER — "This is the core of what I want to share today. A research-backed framework for understanding exactly where AI adoption breaks." """)

# ============================================================
# SLIDE 9: Six Sources — The Framework
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "Behavior Is Shaped by Six Forces — Miss One, Adoption Stalls",
         font_size=24, color=NAVY)
add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
         "The fastest organizations are 12x quicker at turning new ideas into habits. The lever is not tools — it's influence.",
         font_size=13, color=NEAR_BLACK)

# Headers
add_box(slide, Inches(2.6), Inches(1.5), Inches(5), Inches(0.5), NAVY)
add_text(slide, Inches(2.6), Inches(1.5), Inches(5), Inches(0.5),
         "MOTIVATION", font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_box(slide, Inches(7.8), Inches(1.5), Inches(5), Inches(0.5), NAVY)
add_text(slide, Inches(7.8), Inches(1.5), Inches(5), Inches(0.5),
         "ABILITY", font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Row labels
rows = [("PERSONAL", Inches(2.2)), ("SOCIAL", Inches(3.8)), ("STRUCTURAL", Inches(5.4))]
for label, y in rows:
    add_box(slide, Inches(0.6), y, Inches(1.8), Inches(1.4), NAVY)
    add_text(slide, Inches(0.6), y + Inches(0.4), Inches(1.8), Inches(0.5),
             label, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# 6 source cells
sources = [
    # Row 1: Personal
    (Inches(2.6), Inches(2.2), "Source 1: Personal Motivation",
     "Do people WANT to use AI?\nFear of replacement, loss of identity, distrust of algorithms"),
    (Inches(7.8), Inches(2.2), "Source 2: Personal Ability",
     "Do people KNOW HOW to use AI?\nPrompting skills, output validation, workflow integration"),
    # Row 2: Social
    (Inches(2.6), Inches(3.8), "Source 3: Social Motivation",
     "Do peers and leaders MODEL AI use?\nOr is AI adoption seen as optional / risky?"),
    (Inches(7.8), Inches(3.8), "Source 4: Social Ability",
     "Do teams SUPPORT each other in adoption?\nOr is everyone figuring it out alone?"),
    # Row 3: Structural
    (Inches(2.6), Inches(5.4), "Source 5: Structural Motivation",
     "Do INCENTIVES reward AI adoption?\nOr do they reward the old way of working?"),
    (Inches(7.8), Inches(5.4), "Source 6: Structural Ability",
     "Does the ENVIRONMENT make AI easy?\nOr is AI bolted on outside core workflows?"),
]

for x, y, title, desc in sources:
    add_box(slide, x, y, Inches(5), Inches(1.4), LAVENDER)
    add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(4.6), Inches(0.4),
             title, font_size=11, bold=True, color=NAVY)
    add_text(slide, x + Inches(0.2), y + Inches(0.5), Inches(4.6), Inches(0.8),
             desc, font_size=10, color=NEAR_BLACK)

add_footer(slide, 6)

add_notes(slide, """SPEAKING POINTS — The Six Sources Framework

THIS IS THE KEY SLIDE. Spend 60-90 seconds here.

INTRODUCE THE GRID:
- "Behavior — including whether people adopt AI — is not driven by a single factor. It's shaped by six forces, organized in a 2x3 grid: personal, social, and structural — each with a motivation side and an ability side."

WALK THE GRID (top to bottom):
- PERSONAL: "Do individuals WANT to use AI? And do they KNOW HOW? Fear of replacement kills motivation. Lack of skills gets disguised as skepticism — 'AI isn't reliable' often means 'I don't know how to prompt it.'"
- SOCIAL: "Do the people around them — peers, leaders — model AI use? Or is adoption seen as optional, even risky? And when someone tries, does the team support them or leave them to figure it out alone?"
- STRUCTURAL: "Do incentives reward AI adoption? Or do they reward the old way? And does the environment make AI easy — embedded in workflows — or bolted on as an extra step?"

THE INSIGHT:
- "If even ONE of these six is misaligned, adoption stalls. If several are misaligned, it fails entirely. This is why 'just train people' doesn't work — training only addresses Source 2 (Personal Ability). You still have five other sources pulling against you."
- "Research across 1,700 professionals found the fastest-adapting organizations are 12x quicker at turning ideas into habits. Not because of better tools. Because they deliberately align all six sources."

TRANSITION:
- "Now let me show you how cultural debt shows up in each of these six sources — and where to look for it."
""")

# ============================================================
# SLIDE 10: Divider — Diagnosing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
navy_bg(slide)
add_text(slide, Inches(0.8), Inches(2.5), Inches(4), Inches(0.8), "04", font_size=48, color=WHITE)
add_text(slide, Inches(0.8), Inches(3.5), Inches(10), Inches(1),
         "Diagnosing Cultural Debt", font_size=40, color=WHITE)
add_text(slide, Inches(0.8), Inches(4.5), Inches(10), Inches(0.6),
         "Through the Six Sources lens", font_size=20, italic=True,
         color=RGBColor(0xBB, 0xBB, 0xFF))

add_notes(slide, """SECTION DIVIDER — "Now we apply the framework to our situation. Where is cultural debt hiding in each of the six sources?" """)

# ============================================================
# SLIDE 11: Cultural Debt Mapped to Six Sources
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "Cultural Debt Shows Up in Every Source — Here's Where to Look",
         font_size=24, color=NAVY)

data = [
    ["Source", "The Cultural Debt", "What It Looks Like", "Signal"],
    ["1. Personal\nMotivation",
     "Fear and distrust",
     "\"AI will replace my role\" — silent resistance, no one experiments",
     "Low voluntary usage despite tool access"],
    ["2. Personal\nAbility",
     "Skill gaps disguised\nas skepticism",
     "\"AI isn't reliable enough\" — actually means \"I don't know how to use it\"",
     "Inconsistent quality across teams"],
    ["3. Social\nMotivation",
     "No visible role models",
     "No leader publicly using AI. Adoption seen as optional or risky.",
     "\"Promising but premature\" verdicts"],
    ["4. Social\nAbility",
     "No peer support\nstructure",
     "Everyone figures it out alone. No shared practices, no councils.",
     "Isolated pockets of use, no scaling"],
    ["5. Structural\nMotivation",
     "Incentives reward\nthe old way",
     "Promotions based on old metrics. No reward for AI adoption.",
     "Rational sabotage by process owners"],
    ["6. Structural\nAbility",
     "AI bolted on,\nnot embedded",
     "Tools exist outside core workflows. Using AI requires extra steps.",
     "High license cost, low actual usage"],
]
add_table(slide, Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.0), 7, 4, data,
          col_widths=[Inches(1.8), Inches(2.5), Inches(4.5), Inches(3.3)], font_size=11)

add_box(slide, Inches(0.6), Inches(6.3), Inches(12), Inches(0.5), LAVENDER)
add_text(slide, Inches(1), Inches(6.35), Inches(11), Inches(0.4),
         "If even one source is misaligned, adoption stalls. If several are misaligned, it fails entirely.",
         font_size=13, bold=True, color=NAVY)

add_footer(slide, 7)

add_notes(slide, """SPEAKING POINTS — Diagnosing Cultural Debt Through Six Sources

THIS IS THE DIAGNOSTIC SLIDE. Walk down each row and make it concrete.

HOW TO DELIVER:
- "Let me walk through each source and show you what cultural debt looks like in practice."

SOURCE 1 — Personal Motivation:
- "Fear and distrust. People quietly worry AI will replace parts of their role. Nobody says it out loud, but nobody experiments either. The signal: low voluntary usage despite having tool access. They have the tools. They choose not to use them."

SOURCE 2 — Personal Ability:
- "Skill gaps disguised as skepticism. When someone says 'AI isn't reliable enough for our work,' they often mean 'I tried it once, got a bad result, and didn't know how to fix it.' The signal: inconsistent quality across teams — some get great results, others dismiss the tools entirely."

SOURCE 3 — Social Motivation:
- "No visible role models. If no leader in the room is publicly using AI in their actual decision-making, the implicit message is: this is optional. The signal: you hear 'promising but premature' in pilot reviews."

SOURCE 4 — Social Ability:
- "No peer support. Everyone's figuring it out alone. There's no shared practice, no community, no place to ask 'how did you prompt that?' The signal: isolated pockets of adoption that never scale."

SOURCE 5 — Structural Motivation:
- "This is the big one. Incentives reward the old way of working. Promotions are based on metrics that predate AI. Process owners whose authority depends on the current workflow have every rational reason to slow adoption. This is Palmer's 'bureaucratic antibodies.' The signal: rational sabotage disguised as diligence."

SOURCE 6 — Structural Ability:
- "AI is bolted on, not embedded. Using it requires leaving your core workflow, opening another tool, copying results back. Every extra step is friction. The signal: high license costs, low actual usage."

PUNCHLINE:
- Point to the bottom bar: "If even one source is misaligned, adoption stalls. Look at this table — how many are misaligned in your organization right now?"
- PAUSE. This is a discussion prompt. If time allows, take the temperature.

TRANSITION:
- "The good news: each of these has a research-proven fix. Let me show you."
""")

# ============================================================
# SLIDE 12: Divider — Unwinding It
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
navy_bg(slide)
add_text(slide, Inches(0.8), Inches(2.5), Inches(4), Inches(0.8), "05", font_size=48, color=WHITE)
add_text(slide, Inches(0.8), Inches(3.5), Inches(10), Inches(1),
         "Unwinding Cultural Debt", font_size=40, color=WHITE)
add_text(slide, Inches(0.8), Inches(4.5), Inches(10), Inches(0.6),
         "Source by source", font_size=20, italic=True,
         color=RGBColor(0xBB, 0xBB, 0xFF))

add_notes(slide, """SECTION DIVIDER — "We've diagnosed the problem across all six sources. Now here's what to do about it — source by source, with real examples." """)

# ============================================================
# SLIDE 13: The Fix — Source by Source
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "Align All Six Sources to Unwind Cultural Debt and Drive AI Adoption",
         font_size=22, color=NAVY)
add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
         "Leaders who engage all six sources increase their odds of rapid change 10x.",
         font_size=13, bold=True, color=MED_BLUE)

# Headers
add_box(slide, Inches(2.6), Inches(1.4), Inches(5), Inches(0.5), NAVY)
add_text(slide, Inches(2.6), Inches(1.4), Inches(5), Inches(0.5),
         "MOTIVATION — The Fix", font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_box(slide, Inches(7.8), Inches(1.4), Inches(5), Inches(0.5), NAVY)
add_text(slide, Inches(7.8), Inches(1.4), Inches(5), Inches(0.5),
         "ABILITY — The Fix", font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

row_labels = [("PERSONAL", Inches(2.1)), ("SOCIAL", Inches(3.7)), ("STRUCTURAL", Inches(5.3))]
for label, y in row_labels:
    add_box(slide, Inches(0.6), y, Inches(1.8), Inches(1.4), NAVY)
    add_text(slide, Inches(0.6), y + Inches(0.4), Inches(1.8), Inches(0.5),
             label, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

fixes = [
    # Personal
    (Inches(2.6), Inches(2.1),
     "Reframe through storytelling",
     "Share stories where AI = strategic advantage, not threat. Energy firm: AI anomaly detection prevented $M in downtime — reframed from \"tool\" to \"safeguard.\""),
    (Inches(7.8), Inches(2.1),
     "Build skills through practice",
     "Hands-on workshops with real scenarios — prompting, validating, integrating. Services firm: confidence and adoption up in one quarter."),
    # Social
    (Inches(2.6), Inches(3.7),
     "Leaders must visibly model AI use",
     "Respected leaders publicly demonstrate AI in decisions. Retail: regional leaders using AI forecasts shifted AI from \"optional\" to \"how high performers operate.\""),
    (Inches(7.8), Inches(3.7),
     "Create cross-functional AI councils",
     "Compliance + business heads review AI applications together. Mfg org: reduced fear, increased trust, accelerated responsible adoption."),
    # Structural
    (Inches(2.6), Inches(5.3),
     "Tie adoption to promotion criteria",
     "Financial services firm: incorporated responsible AI use into promotion criteria. Clearly communicated adoption is a leadership expectation."),
    (Inches(7.8), Inches(5.3),
     "Embed AI into core workflows",
     "Healthcare: AI risk alerts embedded in EMR workflow, requiring review before chart close. Adoption became automatic, not optional."),
]

for x, y, title, desc in fixes:
    add_box(slide, x, y, Inches(5), Inches(1.4), RGBColor(0xE8, 0xF5, 0xE9))
    add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(4.6), Inches(0.35),
             title, font_size=11, bold=True, color=RGBColor(0x1B, 0x5E, 0x20))
    add_text(slide, x + Inches(0.2), y + Inches(0.45), Inches(4.6), Inches(0.85),
             desc, font_size=10, color=NEAR_BLACK)

add_footer(slide, 8)

add_notes(slide, """SPEAKING POINTS — The Fix: Source by Source

THIS IS THE SOLUTION SLIDE. Same 2x3 grid, but now in green — showing the proven fix for each source. Each cell has a real-world example from the Crucial Learning research.

DELIVERY — Walk the grid, matching diagnosis to prescription:

PERSONAL MOTIVATION (Source 1):
- "The fix for fear is not logic — it's storytelling. An energy company shared stories in executive meetings about how AI anomaly detection prevented millions in equipment downtime. That reframed AI from 'experimental tool' to 'strategic safeguard.' Stories change identity. Spreadsheets don't."

PERSONAL ABILITY (Source 2):
- "A global services firm found consultants using AI inconsistently. Instead of mandates, they ran hands-on workshops where teams practiced prompting, validating, and integrating AI into live client scenarios. Adoption and confidence increased measurably within one quarter. Practice beats training."

SOCIAL MOTIVATION (Source 3):
- "At a large retailer, respected regional leaders publicly demonstrated how they used AI-generated demand forecasts. Their visible modeling shifted AI from 'optional innovation' to 'how high performers operate.' If the leaders you admire use it, it stops being scary."

SOCIAL ABILITY (Source 4):
- "A manufacturing org created cross-functional 'AI councils' where compliance leaders and business heads reviewed AI applications together. This reduced fear, increased trust, and accelerated responsible adoption. People need a safe space to learn together."

STRUCTURAL MOTIVATION (Source 5):
- "A financial services firm incorporated responsible AI use into promotion criteria. Not just experimentation — actual adoption. This clearly communicated that adoption is a leadership expectation, not a nice-to-have. You get what you incentivize."

STRUCTURAL ABILITY (Source 6):
- "A healthcare system embedded AI risk alerts directly into their EMR workflow, requiring documentation of review before closing patient charts. Adoption became automatic, not optional. The best behavior design removes the choice architecture entirely."

KEY INSIGHT:
- "Leaders who engage all six sources are not incrementally better — they are 10x more likely to drive rapid change. That's not a marginal improvement. That's a different category of outcome."

TRANSITION:
- "Before we get to actions, one historical parallel that puts all of this in perspective."
""")

# ============================================================
# SLIDE 14: The Electrification Parallel
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "We've Seen This Before — The 40-Year Reorganization Gap",
         font_size=24, color=NAVY)

# Timeline
add_box(slide, Inches(1), Inches(1.5), Inches(4.5), Inches(2.5), LAVENDER)
add_text(slide, Inches(1.3), Inches(1.6), Inches(4), Inches(0.4),
         "Electrification", font_size=18, bold=True, color=NAVY)
add_text(slide, Inches(1.3), Inches(2.1), Inches(4), Inches(1.5),
         "Demonstrated: 1880\nFactory productivity gains: 1920s\n\nThe lag was not technological. Factories had to be physically rebuilt. Foremen had to be retrained or replaced.",
         font_size=12, color=NEAR_BLACK)

add_box(slide, Inches(6), Inches(1.5), Inches(4.5), Inches(2.5), LAVENDER)
add_text(slide, Inches(6.3), Inches(1.6), Inches(4), Inches(0.4),
         "AI Adoption", font_size=18, bold=True, color=NAVY)
add_text(slide, Inches(6.3), Inches(2.1), Inches(4), Inches(1.5),
         "Enterprise AI: 2023-2025\nScaled value realization: TBD\n\nThe reorganization is cultural, not architectural. The \"foremen\" are your direct reports.",
         font_size=12, color=NEAR_BLACK)

# Arrow between
add_text(slide, Inches(5.5), Inches(2.3), Inches(0.5), Inches(0.5),
         "=", font_size=36, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)

# Jevons
add_box(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(1.2), NAVY)
add_text(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.4),
         "Jevons' Paradox: Cheap intelligence won't reduce work. It will create more of it.",
         font_size=16, bold=True, color=WHITE)
add_text(slide, Inches(1), Inches(5.1), Inches(11), Inches(0.5),
         "Cheap coal didn't reduce coal consumption — it democratized energy for the industrial revolution. Cheap intelligence will do the same for work. The question is whether your organization's culture will let it.",
         font_size=12, color=RGBColor(0xCC, 0xCC, 0xFF))

add_text(slide, Inches(0.6), Inches(6.2), Inches(12), Inches(0.3),
         "Source: Shelly Palmer, \"Unwinding Cultural Debt\" (2026)",
         font_size=9, italic=True, color=LIGHT_GRAY)
add_footer(slide, 9)

add_notes(slide, """SPEAKING POINTS — The Electrification Parallel

THIS SLIDE PROVIDES HISTORICAL CONTEXT. Use it to normalize the challenge and reframe urgency.

ELECTRIFICATION:
- "Electric power was demonstrated in 1880. But it didn't show up in factory productivity statistics until the 1920s — a 40-year gap. The lag was not technological. Factories had to be physically rebuilt around electric motors instead of steam-driven belt systems. And the foremen who understood the old layout had to be retrained or replaced. Historians call this the 'reorganization gap.'"
- "We are in an identical gap right now. Except the reorganization is cultural, not architectural. And the 'foremen' are your direct reports."

JEVONS' PARADOX:
- "Here's the optimistic frame. Jevons' Paradox tells us that when something becomes dramatically cheaper, you don't use less of it — you use catastrophically more. Cheap coal didn't reduce coal consumption. It democratized the energy that powered the industrial revolution."
- "Cheap intelligence is unlikely to reduce work. It will create entirely new categories of work — new industries, new capabilities, new value. But only for organizations whose culture allows them to absorb it."

THE PUNCHLINE:
- "The question is not whether AI is powerful enough. The question is whether our organizational culture will let us use it."

TRANSITION:
- "So here are three concrete actions we can take this quarter to start unwinding our cultural debt."
""")

# ============================================================
# SLIDE 15: Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "Three Actions This Quarter — Mapped to the Six Sources",
         font_size=26, color=NAVY)

actions = [
    ("1. Cultural Debt Audit (Sources 5 & 6)",
     "Add a cultural debt audit to every AI-First POD kickoff. Map processes, identify antibodies, ask: \"Would this step exist if we built from scratch?\"",
     "Owner: Architecture + POD Leads  |  By: End of quarter"),
    ("2. Incentive Redesign (Sources 3 & 5)",
     "Tie responsible AI adoption — not just experimentation — to performance reviews and promotion criteria. Make leaders visibly model AI use.",
     "Owner: HR + Transformation Steering Committee  |  By: Next performance cycle"),
    ("3. \"Future-Back\" Workflow Session (Sources 1 & 2)",
     "Run a design session where one POD designs assuming AI is the primary producer and humans evaluate output quality. Build skills through practice.",
     "Owner: AI Workflow Orchestrators  |  By: Within 30 days"),
]

for i, (title, desc, owner) in enumerate(actions):
    y = Inches(1.2) + Inches(i * 1.6)
    color = LAVENDER if i % 2 == 0 else SOFT_LAVENDER
    add_box(slide, Inches(0.6), y, Inches(12), Inches(1.4), color)
    add_text(slide, Inches(1), y + Inches(0.1), Inches(11), Inches(0.4),
             title, font_size=15, bold=True, color=NAVY)
    add_text(slide, Inches(1), y + Inches(0.55), Inches(11), Inches(0.4),
             desc, font_size=12, color=NEAR_BLACK)
    add_text(slide, Inches(1), y + Inches(1.0), Inches(11), Inches(0.3),
             owner, font_size=10, bold=True, color=MED_BLUE)

# Closing
add_box(slide, Inches(0.6), Inches(6.1), Inches(12), Inches(0.6), NAVY)
add_text(slide, Inches(1), Inches(6.15), Inches(11), Inches(0.5),
         "The question is not whether you've invested in AI. The question is whether you've invested in the human system required to make it work.",
         font_size=13, bold=True, color=WHITE, italic=True)

add_footer(slide, 10)

add_notes(slide, """SPEAKING POINTS — Three Actions This Quarter

THIS IS YOUR CLOSING CONTENT SLIDE. These are the concrete asks. Deliver with conviction.

FRAMING:
- "I'm not asking for a new program or a new budget. I'm asking for three targeted actions, each mapped to specific sources of influence we just diagnosed."

ACTION 1 — Cultural Debt Audit (Sources 5 & 6):
- "First, add a cultural debt audit to every AI-First POD kickoff. Before we encode any process into an AI agent, we need to ask: would this step exist if we built from scratch today? Map the processes AI will touch. Identify the antibodies — the roles whose authority depends on current workflows. This is cheap and fast, and it prevents us from automating dysfunction."
- "This addresses structural motivation and structural ability — the incentives and environment that block adoption."

ACTION 2 — Incentive Redesign (Sources 3 & 5):
- "Second, tie responsible AI adoption — not just experimentation — to performance reviews and promotion criteria. Right now, the rational career move is to hedge. We need to make adoption the rational career move."
- "This addresses social motivation (leaders modeling) and structural motivation (what gets rewarded)."

ACTION 3 — Future-Back Workflow (Sources 1 & 2):
- "Third, run one 'future-back' design session where a POD designs their ideal workflow assuming AI is the primary producer and humans evaluate quality. This isn't a pilot — it's a thought exercise that builds personal motivation and personal ability simultaneously."

THE CLOSING LINE:
- Read the blue bar aloud: "The question is not whether you've invested in AI. The question is whether you've invested in the human system required to make it work."
- Pause. Let it land. This is the line they'll remember.

IF TIME ALLOWS:
- Open for discussion: "Which of these three feels most urgent for us? Where do you see the most cultural debt?"
""")

# ============================================================
# SLIDE 16: Appendix Divider
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
navy_bg(slide)
add_text(slide, Inches(0.8), Inches(3), Inches(10), Inches(1.2),
         "Appendix", font_size=44, color=WHITE)

add_notes(slide, """APPENDIX DIVIDER

- If you are within time, say: "That's the core presentation. What follows is backup material — anticipated questions, full statistics, and references. Happy to go through any of it."
- If you are over time, skip directly to Q&A from the audience rather than walking through appendix slides.
- These slides exist so you can flip to them if a specific question arises during discussion.
""")

# ============================================================
# SLIDE 17: Q&A
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "Anticipated Questions", font_size=28, color=NAVY)

questions = [
    ("\"Isn't this just change management?\"",
     "Traditional change management covers 2 of 6 sources (training + communication). That's why it has a 70% failure rate. This framework covers all six."),
    ("\"How do we measure cultural debt?\"",
     "Three proxies: (1) Time from pilot to production — long = antibodies. (2) Process steps no one can explain. (3) Gap between \"AI deployed\" vs. \"AI used weekly.\""),
    ("\"Won't this slow down AI-First delivery?\"",
     "Cultural debt IS the bottleneck right now. Every \"promising but premature\" verdict costs more than an audit ever will."),
    ("\"Are middle managers the problem?\"",
     "No. Their incentives are (Source 5). Redesign the incentives and the antibody response disappears. It's structural, not personal."),
    ("\"How does this connect to our AI-First Operating Model?\"",
     "The operating model defines WHAT we build (PODs, hub-spoke, Human + AI roles). Six Sources defines WHETHER it gets adopted."),
]

for i, (q, a) in enumerate(questions):
    y = Inches(1.0) + Inches(i * 1.1)
    add_box(slide, Inches(0.6), y, Inches(5), Inches(0.9), LAVENDER if i % 2 == 0 else SOFT_LAVENDER)
    add_text(slide, Inches(0.8), y + Inches(0.15), Inches(4.6), Inches(0.6),
             q, font_size=12, bold=True, color=NAVY)
    add_text(slide, Inches(5.8), y + Inches(0.05), Inches(7), Inches(0.8),
             a, font_size=11, color=NEAR_BLACK)

add_footer(slide, 11)

add_notes(slide, """SPEAKING POINTS — Anticipated Questions

PURPOSE: This slide is your safety net. You don't present it — you flip to it when someone asks one of these questions.

DELIVERY GUIDANCE FOR EACH:

Q1 — "Isn't this just change management?"
- This is the most common objection. Be ready for it.
- Key rebuttal: "Traditional change management focuses on communication and training — that's Sources 1 and 2 out of six. It's like treating a six-cylinder engine with two spark plugs. The other four sources — social motivation, social ability, structural motivation, structural ability — are where most adoption actually stalls. That's why change management has a 70% failure rate."

Q2 — "How do we measure cultural debt?"
- Leaders love metrics. Give them three concrete proxies they can check this week:
  (1) Time from pilot approval to first production use — if it's long, antibodies are at work.
  (2) Process steps that exist but no one can explain why — map them, you'll find cultural debt.
  (3) The gap between 'AI tools deployed' and 'AI tools actually used weekly' — that delta IS your cultural debt, quantified.

Q3 — "Won't this slow down delivery?"
- Counterintuitive but true: "Cultural debt is the thing slowing delivery right now. Every stalled pilot, every 'promising but premature' verdict, every re-evaluation costs more time and budget than a cultural audit ever will. This is an investment that accelerates, not decelerates."

Q4 — "Are middle managers the problem?"
- Be careful here — this can feel like blaming. Reframe immediately:
- "No. Middle managers are responding rationally to their incentives. Source 5 — Structural Motivation — is the problem, not the people. If your promotion criteria reward the old way of working, rational actors will preserve the old way. Redesign the incentives and the antibody response disappears."

Q5 — "How does this connect to our AI-First Operating Model?"
- "The operating model defines WHAT we build — PODs, hub-and-spoke, Human + AI roles. The Six Sources framework defines WHETHER it actually gets adopted. One without the other produces investment without return."
""")

# ============================================================
# SLIDE 18: Key Statistics
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "Key Statistics Reference", font_size=28, color=NAVY)

data = [
    ["Statistic", "Source", "Year"],
    ["Only 1 in 4 AI initiatives deliver expected ROI", "IBM Global AI Survey / Fortune", "2025"],
    ["95% of AI initiatives fail to deliver value", "HBR — De Cremer, Schweitzer et al.", "2025"],
    ["90% of orgs use AI in at least one function", "McKinsey State of AI Global Survey", "2025"],
    ["82% of leaders use GenAI weekly", "Wharton AI Adoption Report", "2025"],
    ["12x faster adoption in agile organizations", "Crucial Learning (1,700+ professionals)", "2026"],
    ["10x higher odds with six-source approach", "MIT Sloan Management Review", "2008"],
    ["Electrification: 40-year reorganization gap", "Economic historians", "Historical"],
    ["70% of change management efforts fail", "McKinsey & Company", "2015"],
]
add_table(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(4.5), 9, 3, data,
          col_widths=[Inches(5.5), Inches(4.5), Inches(2)], font_size=12)

add_footer(slide, 12)

# ============================================================
# SLIDE 19: References
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "References", font_size=28, color=NAVY)

refs = [
    ("1", "Fortune (via Yahoo Finance). \"CEOs say only about one in four AI initiatives have delivered expected ROI\" (2025).",
     "finance.yahoo.com/news/ceos-just-fraction-ai-initiatives-173518010.html"),
    ("2", "The Wall Street Journal. \"Companies are struggling to drive a return on AI investment\" (2025).",
     "wsj.com/articles/companies-are-struggling-to-drive-a-return-on-ai"),
    ("3", "McKinsey & Company. \"The State of AI: Global Survey 2025.\"",
     "mckinsey.com/capabilities/quantumblack/our-insights/the-state-of-ai"),
    ("4", "Knowledge at Wharton. \"2025 AI Adoption Report: Gen AI Fast-Tracks Into the Enterprise.\"",
     "knowledge.wharton.upenn.edu/special-report/2025-ai-adoption-report/"),
    ("5", "OpenAI. \"The State of Enterprise AI 2025.\"",
     "openai.com/index/the-state-of-enterprise-ai-2025-report/"),
    ("6", "De Cremer, Schweitzer, McGuire & Narayanan. \"How Behavioral Science Can Improve the Return on AI Investments.\" Harvard Business Review, Nov 2025.",
     ""),
    ("7", "Grenny, Maxfield & Shimberg. \"How to Have Influence.\" MIT Sloan Management Review, Oct 2008.",
     ""),
    ("8", "Crucial Learning. \"Why AI Initiatives Fail: Ignoring Your Human System is the Fastest Path to Wasted AI Investment\" (2026).",
     "cruciallearning.com"),
    ("9", "Shelly Palmer. \"Unwinding Cultural Debt\" (Apr 2026).",
     "shellypalmer.com/2026/04/unwinding-cultural-debt/"),
]

for i, (num, citation, url) in enumerate(refs):
    y = Inches(1.0) + Inches(i * 0.62)
    add_text(slide, Inches(0.6), y, Inches(0.5), Inches(0.5),
             num, font_size=11, bold=True, color=NAVY)
    add_text(slide, Inches(1.1), y, Inches(7.5), Inches(0.55),
             citation, font_size=10, color=NEAR_BLACK)
    if url:
        add_text(slide, Inches(8.8), y + Inches(0.05), Inches(4), Inches(0.45),
                 url, font_size=8, color=MED_BLUE, italic=True)

add_footer(slide, 13)

# ============================================================
# SAVE
# ============================================================
output = "/Users/GUNDLLX/learn-claude/presentations/why-ai-initiatives-fail-editable.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Slides: {len(prs.slides)}")
