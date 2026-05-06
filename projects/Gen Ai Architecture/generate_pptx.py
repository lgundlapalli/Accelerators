"""
Abbott GenAI Blueprint — PowerPoint Generator
"""

import importlib.util
from pathlib import Path

# Load Abbott brand constants
spec = importlib.util.spec_from_file_location(
    "abc", Path.home() / '.claude/commands/pptx-generator/references/abbott-brand-constants.py')
abc = importlib.util.module_from_spec(spec)
spec.loader.exec_module(abc)
for k, v in vars(abc).items():
    if not k.startswith('_'):
        globals()[k] = v

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUT = "/Users/GUNDLLX/learn-claude/projects/Gen Ai Architecture/Abbott-GenAI-Blueprint-Capability-Model-2026-05-06.pptx"

DOMAIN_DATA = [
    {
        "name": "Content & Creative Studio",
        "definition": "Generation, drafting, and transformation of written, visual, and multimedia content for internal and external audiences.",
        "subs": ["Creative Content Draft", "Writing Drafts & Data Entry", "Personalize Content", "Translation"],
        "tools": "Writer Newsroom, Adobe Firefly, AutogenAI",
        "divisions": "ADC, Corp, EPD, GMEA, RMDx, CHR, MD (7 divisions)",
        "maturity": "Scaling",
        "use_cases": 52,
    },
    {
        "name": "Insights & Analytics Acceleration",
        "definition": "AI-powered search, synthesis, and generation of insights from structured and unstructured data.",
        "subs": ["Search, Screen & Insights", "Reporting & Observability", "Forecast"],
        "tools": "Mktg Insights Accelerator, RADIA Research, Databricks",
        "divisions": "GMEA, AN, RMDx, BTS, ADC, Corp (6 divisions)",
        "maturity": "Producing",
        "use_cases": 48,
    },
    {
        "name": "Productivity & Automation",
        "definition": "Enterprise-wide productivity enhancement through AI-integrated tools — meeting summarization, workflow automation, and process intelligence.",
        "subs": ["Enterprise Productivity Tools", "Writing Drafts & Data Entry"],
        "tools": "M365 Copilot, Cursor/Windsurf, Five9 GSD",
        "divisions": "BTS, Corp, CHR, AN, MD, EPD (6 divisions)",
        "maturity": "Scaling",
        "use_cases": 42,
    },
    {
        "name": "Customer & Field Engagement",
        "definition": "AI-powered tools that augment customer-facing roles — field sales, customer service, and HCP interactions.",
        "subs": ["Assistants", "Personalize Content", "Search & CRM Insights"],
        "tools": "SmartRep Call Planning, Neuron7, Five9 GSD",
        "divisions": "EPD, ADC, CoreDx, BTS, GMEA (5 divisions)",
        "maturity": "Producing",
        "use_cases": 35,
    },
    {
        "name": "Engineering & Development",
        "definition": "AI tools augmenting software engineering, quality assurance, and product development workflows.",
        "subs": ["Code Creation with Platform", "Code Creation IDE & Test Authoring"],
        "tools": "Cursor/Windsurf/Atlassian AI, Azure AI Sandbox",
        "divisions": "BTS, MD, RMDx (3 divisions)",
        "maturity": "Scaling",
        "use_cases": 28,
    },
    {
        "name": "Learning, Compliance & Governance",
        "definition": "AI applications supporting employee learning, regulatory compliance, risk management, and organizational knowledge management.",
        "subs": ["Learning, Training & Coaching", "Image Identification / Analysis"],
        "tools": "Pilots in CHR, Corp, RMDx, Cyber",
        "divisions": "CHR, Corp, RMDx, MD, AN, Cyber (6 divisions)",
        "maturity": "Exploring",
        "use_cases": 15,
    },
]

DOMAIN_COLORS_HEX = [
    (0x44, 0x72, 0xC4),
    (0xED, 0x7D, 0x31),
    (0xA9, 0xD1, 0x8E),
    (0xFF, 0x99, 0xCC),
    (0xFF, 0xD9, 0x66),
    (0x9D, 0xC3, 0xE6),
]


def domain_color(i):
    r, g, b = DOMAIN_COLORS_HEX[i]
    return RGBColor(r, g, b)


def maturity_color(mat):
    return {
        "Producing": GREEN_ACCENT,
        "Scaling": MED_BLUE,
        "Exploring": AMBER_ACCENT,
    }.get(mat, LIGHT_GRAY)


prs = create_presentation()

# ── Slide 1: Title ───────────────────────────────────────────
s = add_title_slide(
    prs,
    "Abbott GenAI Blueprint\n& Capability Model",
    "Enterprise AI Landscape — 220 Use Cases Across 15 Divisions",
    tagline="Building a scalable, governed GenAI capability foundation for Abbott",
    org_date="Abbott BTS-DTS | GenAI Center of Excellence | May 2026",
)

# ── Slide 2: Agenda ──────────────────────────────────────────
add_agenda_slide(prs, "Agenda", [
    ("01", "The GenAI Landscape at Abbott", "220 use cases, 12 divisions, 1Q26 production snapshot"),
    ("02", "Six Capability Domains", "Logical grouping of 15 raw categories into strategic domains"),
    ("03", "Enterprise Blueprint", "3-tier platform architecture and CoE governance model"),
    ("04", "Platform Foundation", "Enterprise AI platform investments and rationalization"),
    ("05", "Next Steps", "Priority actions for CoE and division leads"),
], 2)

# ── Slide 3: Section Divider — Landscape ────────────────────
add_section_divider(prs, "01", "The GenAI Landscape at Abbott",
                    "1Q26 Portfolio Snapshot — 220 Use Cases")

# ── Slide 4: Big Numbers ─────────────────────────────────────
add_big_numbers_slide(
    prs,
    "Abbott GenAI — By the Numbers",
    "1Q26 snapshot across all divisions and capability categories",
    [
        ("220", "Total Use Cases", "Catalogued across 12 divisions and 15 capability types"),
        ("40", "In Production", "18% of portfolio actively deployed and generating value"),
        ("12", "Divisions", "BTS, ADC, Corp, AN, RMDx, MD, EPD, CHR, GMEA, CoreDx, Lingo, Cyber"),
        ("61%", "Buy vs Build", "133 Buy / 86 Build — platform-first, custom-when-needed"),
    ],
    "43 use cases in active WIP + 29 pending ESC approval — accelerating the pipeline is the priority",
    "Source: Abbott GenAI Historical Catalog of Use Cases — 1Q26",
    4,
)

# ── Slide 5: Section Divider — Domains ──────────────────────
add_section_divider(prs, "02", "Six Capability Domains",
                    "Consolidating 15 categories into 6 strategic capability areas")

# ── Slide 6: Domain Overview Table ──────────────────────────
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide6, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "Six GenAI Capability Domains — Overview", font_size=26, color=NAVY)
headers = ["Domain", "Definition (brief)", "Key Sub-Capabilities", "Divisions", "Maturity"]
rows_data = [headers] + [
    [
        d["name"],
        d["definition"][:80] + "...",
        ", ".join(d["subs"][:2]),
        d["divisions"].split("(")[0].strip(),
        d["maturity"],
    ]
    for d in DOMAIN_DATA
]
add_table(slide6, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.5),
          7, 5, rows_data,
          col_widths=[Inches(2.3), Inches(3.0), Inches(2.6), Inches(2.2), Inches(1.2)],
          font_size=10)
add_footer(slide6, 6)

# ── Slides 7-12: One per domain ──────────────────────────────
for pg, d in enumerate(DOMAIN_DATA, start=7):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    idx = pg - 7
    dc = domain_color(idx)

    # Colored header bar
    add_box(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), dc)
    add_text(slide, Inches(0.6), Inches(0.1), Inches(11), Inches(0.7),
             f"Domain {idx + 1}: {d['name']}",
             font_size=22, bold=True, color=WHITE)

    # Maturity badge
    mc = maturity_color(d["maturity"])
    add_box(slide, Inches(11.3), Inches(0.15), Inches(1.7), Inches(0.55), mc)
    add_text(slide, Inches(11.3), Inches(0.15), Inches(1.7), Inches(0.55),
             d["maturity"], font_size=13, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)

    # Definition
    add_text(slide, Inches(0.6), Inches(1.05), Inches(12), Inches(0.5),
             d["definition"], font_size=12, color=NEAR_BLACK, italic=True)

    # Main table
    sub_rows = [[s, "Active"] for s in d["subs"]]
    headers2 = ["Sub-Capability", "Status"]
    tdata = [headers2] + sub_rows
    add_table(slide, Inches(0.6), Inches(1.7), Inches(6.0), Inches(3.2),
              len(tdata), 2, tdata,
              col_widths=[Inches(4.5), Inches(1.5)],
              font_size=11)

    # Right panel info
    add_box(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(3.2), SOFT_LAVENDER)
    add_text(slide, Inches(7.2), Inches(1.8), Inches(5.4), Inches(0.35),
             "In-Production Tools", font_size=12, bold=True, color=NAVY)
    add_text(slide, Inches(7.2), Inches(2.2), Inches(5.4), Inches(0.8),
             d["tools"], font_size=11, color=NEAR_BLACK)
    add_text(slide, Inches(7.2), Inches(3.1), Inches(5.4), Inches(0.35),
             "Divisions Using", font_size=12, bold=True, color=NAVY)
    add_text(slide, Inches(7.2), Inches(3.5), Inches(5.4), Inches(0.6),
             d["divisions"], font_size=11, color=NEAR_BLACK)
    add_text(slide, Inches(7.2), Inches(4.2), Inches(5.4), Inches(0.35),
             f"Use Cases: {d['use_cases']} (estimated)", font_size=11,
             bold=True, color=MED_BLUE)

    add_footer(slide, pg)

# ── Slide 13: Section Divider — Blueprint ───────────────────
add_section_divider(prs, "03", "Enterprise Blueprint",
                    "3-Tier Architecture — Platform, Capabilities, Business Domains")

# ── Slide 14: Blueprint Architecture ────────────────────────
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide14, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "Abbott GenAI Enterprise Blueprint — 3-Tier Architecture", font_size=24, color=NAVY)

tiers = [
    ["TIER 1\nBUSINESS DOMAINS",
     "Content & Creative Studio | Insights & Analytics | Productivity & Automation",
     "Customer & Field Engagement | Engineering & Development | Learning & Governance"],
    ["TIER 2\nGENAI CAPABILITIES",
     "Creative Content Draft & Translation | Search & Insights | Enterprise Productivity",
     "Assistants & Personalization | Code Creation | Learning & Compliance"],
    ["TIER 3\nSHARED PLATFORM",
     "Azure AI Foundry | Anthropic Claude | Salesforce Einstein",
     "Adobe Experience Cloud | Veeva Vault AI | Microsoft M365 Copilot"],
]
tier_colors = [LAVENDER, SOFT_LAVENDER, RGBColor(0xE8, 0xF0, 0xFE)]
for ti, (tier, tc) in enumerate(zip(tiers, tier_colors)):
    y = Inches(1.1) + Inches(ti * 1.7)
    add_box(slide14, Inches(0.5), y, Inches(12.3), Inches(1.55), tc)
    add_box(slide14, Inches(0.5), y, Inches(1.8), Inches(1.55), NAVY)
    add_text(slide14, Inches(0.55), y + Inches(0.1), Inches(1.7), Inches(1.35),
             tier[0], font_size=9, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text(slide14, Inches(2.5), y + Inches(0.1), Inches(10.1), Inches(0.55),
             tier[1], font_size=11, bold=True, color=NAVY)
    add_text(slide14, Inches(2.5), y + Inches(0.7), Inches(10.1), Inches(0.65),
             tier[2], font_size=11, color=NEAR_BLACK)

# Arrows
for ai in range(2):
    y_arrow = Inches(2.65) + Inches(ai * 1.7)
    add_box(slide14, Inches(6.4), y_arrow, Inches(0.5), Inches(0.2), LIGHT_GRAY)

add_text(slide14, Inches(0.5), Inches(6.1), Inches(12), Inches(0.35),
         "CoE Governance: ESC review, platform standards, security guardrails, and division enablement",
         font_size=10, italic=True, color=LIGHT_GRAY)
add_footer(slide14, 14)

# ── Slide 15: Platform Foundation ───────────────────────────
slide15 = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide15, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "Enterprise Platform Foundation", font_size=26, color=NAVY)
plat_data = [
    ["Platform", "Primary Use Cases", "Domains Served", "Deployment"],
    ["Azure AI Foundry", "Model access, AI infrastructure, custom fine-tuning", "All domains", "Enterprise"],
    ["Anthropic Claude", "Long-context reasoning, document analysis, code review", "Insights, Engineering, Compliance", "Enterprise API"],
    ["Salesforce Einstein", "CRM-native AI, sales copilot, next-best-action", "Customer & Field Engagement", "SaaS embedded"],
    ["Adobe Experience Cloud", "Content generation, creative tools, personalization", "Content & Creative Studio", "SaaS"],
    ["Veeva Vault AI", "Regulated content mgmt, MLR review, clinical docs", "Compliance & Governance", "SaaS"],
    ["Microsoft M365 Copilot", "Enterprise productivity, Teams, Office suite AI", "Productivity & Automation", "M365 tenant"],
]
add_table(slide15, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.5),
          7, 4, plat_data,
          col_widths=[Inches(2.5), Inches(4.5), Inches(3.0), Inches(2.3)],
          font_size=10)
add_footer(slide15, 15)

# ── Slide 16: Next Steps ─────────────────────────────────────
add_action_cards_slide(prs, "Next Steps & Priority Actions", [
    ("Consolidate Content Platforms",
     "Rationalize Writer, Adobe Firefly, and AutogenAI under a unified CoE contract",
     "Owners: ADC + GMEA + BTS platform leads | Timeline: Q2 2026",
     "Owner: BTS-DTS CoE + ADC"),
    ("Accelerate ESC Approvals",
     "29 use cases pending — establish 30-day SLA for ESC review cadence",
     "Target: clear backlog by end of Q2 2026 | Impact: ~$Xm deferred value",
     "Owner: ESC Chair + CoE Program Lead"),
    ("Expand SmartRep Model to ADC & GMEA",
     "Template EPD's SmartRep success for replication across commercial divisions",
     "Playbook: data requirements, Salesforce Einstein integration, training",
     "Owner: EPD + ADC Commercial AI Leads"),
    ("Launch Domain Champions Program",
     "Assign a CoE domain champion per capability domain to drive reuse and sharing",
     "6 champions x 6 domains | Quarterly cross-division syncs",
     "Owner: GenAI CoE Director"),
    ("Investment Plan for Learning & Governance Domain",
     "Develop pilot use case in CHR to advance the Exploring → Scaling maturity",
     "Scope: 1 pilot use case, Q3 2026, with compliance and legal review",
     "Owner: CHR + Corp + Cyber"),
],
    "GenAI is a strategic priority — pace of production deployment is the leading indicator of value realization",
    16,
)

prs.save(OUT)
print(f"  Saved {OUT}")
